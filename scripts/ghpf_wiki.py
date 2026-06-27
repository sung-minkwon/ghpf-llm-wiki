#!/usr/bin/env python3
"""Sidecar utilities for GHFP LLM Wiki vaults."""

from __future__ import annotations

import argparse
import hashlib
import html
import json
import re
import shutil
import sqlite3
import subprocess
import tempfile
import urllib.parse
import urllib.request
from datetime import datetime, timedelta, timezone
from difflib import SequenceMatcher
from html.parser import HTMLParser
from pathlib import Path

WIKILINK_RE = re.compile(r"\[\[([^\]|#]+)(?:#[^\]|]+)?(?:\|[^\]]+)?\]\]")
QUALITY_REQUIRED_FIELDS = ["tags", "source", "created", "aliases"]
QUALITY_WEIGHTS = {
    "completeness": 0.20,
    "connections": 0.20,
    "coverage": 0.25,
    "consistency": 0.20,
    "suggestions": 0.15,
}
NON_RETRIEVAL_PAGE_NAMES = {"index.md", "log.md", "overview.md"}
NON_RETRIEVAL_SECTIONS = {"## Related Notes", "## Backlinks", "## Key Links"}
PIPELINE_STEPS = ["setup", "ingest", "compile", "lint", "strengthen", "file-back", "graph", "context"]
GRAPHIFY_RAW_DIR = ("raw", "graphify_articles")
INTERNAL_RAW_PREFIXES = {
    ("raw", "originals"),
    ("raw", "sources"),
    ("raw", "sources", "extracted"),
    ("raw", "sources", "downloads"),
    ("raw", "figures", "video-frames"),
}
TEXT_SOURCE_SUFFIXES = {".md", ".txt", ".csv", ".json", ".yaml", ".yml"}
OFFICE_SOURCE_SUFFIXES = {".docx", ".pptx", ".xlsx", ".xls", ".hwp", ".hwpx", ".hwpml"}
EXTRACTABLE_SUFFIXES = TEXT_SOURCE_SUFFIXES | {".pdf", ".html", ".htm"} | OFFICE_SOURCE_SUFFIXES
IMAGE_SUFFIXES = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tif", ".tiff"}
VIDEO_SUFFIXES = {".mp4", ".mov", ".mkv", ".webm", ".avi", ".m4v"}
URL_RE = re.compile(r"^https?://", re.IGNORECASE)
VECTOR_DIMS = 256
IMAGE_CONTEXT_WORDS = {"chart", "graph", "plot", "figure", "diagram", "architecture", "flow", "screenshot", "canvas", "svg"}
IMAGE_SKIP_WORDS = {"ads", "doubleclick", "googlesyndication", "tracker", "pixel", "icon", "avatar", "logo", "social", "share", "nav", "menu", "header", "footer"}


def slugify(text: str) -> str:
    slug = re.sub(r"[^A-Za-z0-9가-힣._-]+", "-", text.strip()).strip("-").lower()
    return slug or "untitled"


def markdown_files(vault: Path):
    wiki = vault / "wiki"
    if not wiki.exists():
        return []
    return sorted(p for p in wiki.rglob("*.md") if p.is_file())


def reference_markdown_files(vault: Path) -> list[Path]:
    files = list(markdown_files(vault))
    graph_imports = vault / "graph_imports"
    if graph_imports.exists():
        files.extend(sorted(p for p in graph_imports.rglob("*.md") if p.is_file()))
    return files


def is_retrieval_page(path: Path) -> bool:
    return path.name not in NON_RETRIEVAL_PAGE_NAMES


def find_obsidian_vaults(root: Path, max_depth: int = 3) -> list[Path]:
    root = root.expanduser()
    if not root.exists() or not root.is_dir():
        return []
    found = []
    skip_names = {".git", ".obsidian", "__pycache__", "node_modules", ".venv", "venv"}
    stack = [(root, 0)]
    while stack:
        current, depth = stack.pop()
        if (current / ".obsidian").is_dir():
            found.append(current)
        if depth >= max_depth:
            continue
        try:
            children = list(current.iterdir())
        except OSError:
            continue
        for child in children:
            if child.is_dir() and child.name not in skip_names:
                stack.append((child, depth + 1))
    return sorted(set(found))


def load_vault_config(vault: Path) -> dict:
    config_path = vault / "ghpf.config.json"
    if not config_path.exists():
        return {}
    try:
        return json.loads(config_path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return {"_error": "invalid_json"}


def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def read_text(path: Path, max_chars: int = 50000) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")[:max_chars]


def parse_frontmatter(content: str) -> tuple[dict, str]:
    if not content.startswith("---\n"):
        return {}, content
    parts = content.split("\n---\n", 1)
    if len(parts) != 2:
        return {}, content
    frontmatter = {}
    for line in parts[0].splitlines()[1:]:
        if ":" not in line:
            continue
        key, value = line.split(":", 1)
        key = key.strip()
        value = value.strip()
        if value.startswith("[") and value.endswith("]"):
            items = [item.strip().strip("\"'") for item in value[1:-1].split(",") if item.strip()]
            frontmatter[key] = items
        elif value:
            frontmatter[key] = value.strip("\"'")
    return frontmatter, parts[1]


def frontmatter_block(fields: dict) -> str:
    lines = ["---"]
    for key, value in fields.items():
        if isinstance(value, list):
            rendered = ", ".join(json.dumps(str(item), ensure_ascii=False) for item in value)
            lines.append(f"{key}: [{rendered}]")
        else:
            lines.append(f"{key}: {json.dumps(str(value), ensure_ascii=False)}")
    lines.append("---")
    return "\n".join(lines) + "\n\n"


def manifest_path(vault: Path) -> Path:
    return vault / "wiki" / "manifest.json"


def load_manifest(vault: Path) -> dict:
    path = manifest_path(vault)
    if not path.exists():
        return {"sources": [], "generated_pages": [], "operations": []}
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return {"sources": [], "generated_pages": [], "operations": [{"type": "manifest_recovered", "time": now_iso()}]}


def save_manifest(vault: Path, manifest: dict) -> None:
    path = manifest_path(vault)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def sha256_text(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8", errors="ignore")).hexdigest()


def jsonable(value):
    if isinstance(value, Path):
        return value.as_posix()
    if isinstance(value, dict):
        return {key: jsonable(item) for key, item in value.items()}
    if isinstance(value, list):
        return [jsonable(item) for item in value]
    return value


def titleize_slug(slug: str) -> str:
    return re.sub(r"[-_]+", " ", slug).strip().title()


def append_log(vault: Path, message: str) -> None:
    log_path = vault / "wiki" / "log.md"
    log_path.parent.mkdir(parents=True, exist_ok=True)
    if not log_path.exists():
        log_path.write_text("# GHFP LLM Wiki Log\n\n", encoding="utf-8")
    with log_path.open("a", encoding="utf-8") as handle:
        handle.write(f"- {now_iso()}: {message}\n")


def ensure_index_entry(vault: Path, page_rel: str, title: str) -> None:
    index_path = vault / "wiki" / "index.md"
    index_path.parent.mkdir(parents=True, exist_ok=True)
    if not index_path.exists():
        index_path.write_text("# GHFP LLM Wiki Index\n\n## Pages\n\n", encoding="utf-8")
    text = index_path.read_text(encoding="utf-8")
    entry = f"- [[{title}]] - `{page_rel}`"
    if entry not in text:
        with index_path.open("a", encoding="utf-8") as handle:
            if not text.endswith("\n"):
                handle.write("\n")
            handle.write(entry + "\n")


def source_candidates(vault: Path) -> list[Path]:
    candidates = []
    for root_name in ("raw", "_raw"):
        root = vault / root_name
        if not root.exists():
            continue
        for path in root.rglob("*"):
            rel_parts = path.relative_to(vault).parts
            if rel_parts[:2] == GRAPHIFY_RAW_DIR:
                continue
            if any(rel_parts[: len(prefix)] == prefix for prefix in INTERNAL_RAW_PREFIXES):
                continue
            if path.is_file() and path.suffix.lower() in EXTRACTABLE_SUFFIXES:
                candidates.append(path)
    return sorted(set(candidates))


def copy_to_raw(vault: Path, source: Path) -> Path:
    raw_root = vault / "raw" / "sources"
    raw_root.mkdir(parents=True, exist_ok=True)
    if source.resolve().is_relative_to((vault / "raw").resolve()):
        return source
    digest = sha256_file(source)[:12]
    target = raw_root / f"{slugify(source.stem)}-{digest}{source.suffix.lower() or '.txt'}"
    if not target.exists():
        shutil.copy2(source, target)
    return target


def vault_rel(vault: Path, path: Path) -> str:
    try:
        return path.relative_to(vault).as_posix()
    except ValueError:
        return path.as_posix()


def original_asset_path(vault: Path, label: str, digest: str, suffix: str) -> Path:
    root = vault / "raw" / "originals"
    root.mkdir(parents=True, exist_ok=True)
    clean_suffix = suffix.lower() if suffix else ".bin"
    if not clean_suffix.startswith("."):
        clean_suffix = f".{clean_suffix}"
    return root / f"{slugify(Path(label).stem or label)[:80]}-{digest[:12]}{clean_suffix}"


def preserve_original_file(vault: Path, source: Path, source_label: str | None = None) -> dict:
    source = source.expanduser()
    digest = sha256_file(source)
    target = original_asset_path(vault, source_label or source.name, digest, source.suffix or ".bin")
    if not target.exists():
        if source.resolve() == target.resolve():
            pass
        else:
            shutil.copy2(source, target)
    return {
        "kind": "file",
        "source": str(source),
        "path": vault_rel(vault, target),
        "sha256": digest,
        "bytes": target.stat().st_size,
    }


def preserve_original_bytes(vault: Path, label: str, body: bytes, suffix: str, source_url: str | None = None, content_type: str | None = None) -> dict:
    digest = hashlib.sha256(body).hexdigest()
    target = original_asset_path(vault, label, digest, suffix)
    if not target.exists():
        target.write_bytes(body)
    record = {
        "kind": "url" if source_url else "bytes",
        "source": source_url or label,
        "path": vault_rel(vault, target),
        "sha256": digest,
        "bytes": len(body),
    }
    if source_url:
        record["source_url"] = source_url
    if content_type:
        record["content_type"] = content_type
    return record


def external_original_record(source_url: str) -> dict:
    return {"kind": "url", "source": source_url, "source_url": source_url}


def summarize_text(text: str, max_lines: int = 12) -> list[str]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return ["No readable text extracted."]
    return lines[:max_lines]


def extract_terms(text: str, limit: int = 12) -> list[str]:
    words = re.findall(r"[A-Za-z][A-Za-z0-9_-]{3,}|[가-힣]{2,}", text)
    stop = {
        "about",
        "after",
        "also",
        "best",
        "compare",
        "create",
        "data",
        "from",
        "have",
        "into",
        "must",
        "rather",
        "show",
        "should",
        "source",
        "supports",
        "than",
        "that",
        "this",
        "useful",
        "when",
        "wiki",
        "will",
        "with",
    }
    counts = {}
    for word in words:
        key = word.lower()
        if key in stop:
            continue
        if re.fullmatch(r"[a-f0-9]{12,}", key):
            continue
        if len(key) >= 16 and sum(char.isdigit() for char in key) >= 4:
            continue
        counts[key] = counts.get(key, 0) + 1
    return [word for word, _ in sorted(counts.items(), key=lambda item: (-item[1], item[0]))[:limit]]


def source_key_points(text: str, limit: int = 8) -> list[str]:
    lines = summarize_text(text, max_lines=limit)
    points = []
    for line in lines:
        cleaned = re.sub(r"^#+\s*", "", line).strip("-* \t")
        if cleaned:
            points.append(cleaned[:240])
    return points


def page_title_from_path(path: Path) -> str:
    return titleize_slug(slugify(path.stem))


def page_aliases(path: Path) -> set[str]:
    aliases = {path.stem.lower(), titleize_slug(path.stem).lower()}
    content = path.read_text(encoding="utf-8", errors="ignore")
    frontmatter, body = parse_frontmatter(content)
    for alias in frontmatter.get("aliases", []) if isinstance(frontmatter.get("aliases"), list) else []:
        aliases.add(str(alias).lower())
    for line in body.splitlines():
        if line.startswith("# "):
            title = line[2:].strip()
            aliases.add(title.lower())
            aliases.add(title.replace(":", "").lower())
            break
    return aliases


def existing_page_aliases(vault: Path) -> set[str]:
    aliases = set()
    for page in markdown_files(vault):
        aliases.update(page_aliases(page))
    return aliases


def existing_wikilinks_for_terms(vault: Path, terms: list[str], limit: int = 8) -> tuple[list[str], list[str]]:
    aliases = existing_page_aliases(vault)
    links = []
    candidates = []
    for term in terms[:limit]:
        title = titleize_slug(term)
        if title.lower() in aliases:
            links.append(f"[[{title}]]")
        else:
            candidates.append(title)
    return links, candidates


class SimpleHTMLTextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts = []
        self.title_parts = []
        self.skip_depth = 0
        self.in_title = False

    def handle_starttag(self, tag: str, attrs):
        tag = tag.lower()
        if tag in {"script", "style", "noscript", "svg"}:
            self.skip_depth += 1
        if tag == "title":
            self.in_title = True
        if tag in {"p", "div", "section", "article", "header", "footer", "li", "br", "h1", "h2", "h3"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str):
        tag = tag.lower()
        if tag in {"script", "style", "noscript", "svg"} and self.skip_depth:
            self.skip_depth -= 1
        if tag == "title":
            self.in_title = False
        if tag in {"p", "div", "section", "article", "li", "h1", "h2", "h3"}:
            self.parts.append("\n")

    def handle_data(self, data: str):
        cleaned = html.unescape(data).strip()
        if not cleaned:
            return
        if self.in_title:
            self.title_parts.append(cleaned)
        if self.skip_depth == 0:
            self.parts.append(cleaned)

    @property
    def title(self) -> str:
        return " ".join(self.title_parts).strip()

    @property
    def text(self) -> str:
        raw = " ".join(self.parts)
        lines = [re.sub(r"\s+", " ", line).strip() for line in raw.splitlines()]
        return "\n".join(line for line in lines if line)


class ImageCandidateExtractor(HTMLParser):
    def __init__(self, base_url: str | None = None):
        super().__init__()
        self.base_url = base_url
        self.candidates = []
        self.figure_depth = 0
        self.current_figure = None
        self.caption_depth = 0
        self.caption_parts = []

    def handle_starttag(self, tag: str, attrs):
        tag = tag.lower()
        attrs_dict = {key.lower(): value for key, value in attrs if key and value is not None}
        if tag == "figure":
            self.figure_depth += 1
            self.current_figure = {"context": "figure", "caption": ""}
        if tag == "figcaption" and self.figure_depth:
            self.caption_depth += 1
            self.caption_parts = []
        if tag in {"img", "source"}:
            src = attrs_dict.get("src") or attrs_dict.get("data-src") or attrs_dict.get("srcset", "").split(" ")[0]
            if src:
                self._add_candidate(tag, src, attrs_dict)
        if tag in {"canvas", "svg"}:
            self._add_candidate(tag, attrs_dict.get("id") or attrs_dict.get("class") or tag, attrs_dict, virtual=True)

    def handle_endtag(self, tag: str):
        tag = tag.lower()
        if tag == "figcaption" and self.caption_depth:
            self.caption_depth -= 1
            if self.current_figure is not None:
                self.current_figure["caption"] = " ".join(self.caption_parts).strip()
        if tag == "figure" and self.figure_depth:
            self.figure_depth -= 1
            self.current_figure = None

    def handle_data(self, data: str):
        if self.caption_depth:
            cleaned = html.unescape(data).strip()
            if cleaned:
                self.caption_parts.append(cleaned)

    def _add_candidate(self, tag: str, src: str, attrs: dict, virtual: bool = False) -> None:
        context = " ".join(str(attrs.get(key, "")) for key in ("alt", "title", "class", "id", "role")).lower()
        if any(word in context or word in src.lower() for word in IMAGE_SKIP_WORDS):
            return
        width = parse_int(attrs.get("width"))
        height = parse_int(attrs.get("height"))
        if width and height and width < 100 and height < 100:
            return
        reasons = []
        score = 0
        if self.figure_depth:
            score += 3
            reasons.append("figure")
        if tag in {"canvas", "svg"}:
            score += 4
            reasons.append(tag)
        for word in IMAGE_CONTEXT_WORDS:
            if word in context or word in src.lower():
                score += 2
                reasons.append(word)
        if width and height and (width >= 400 or height >= 300):
            score += 2
            reasons.append("large")
        if attrs.get("alt"):
            score += 1
            reasons.append("alt")
        caption = self.current_figure.get("caption", "") if self.current_figure else ""
        candidate = {
            "tag": tag,
            "src": src if virtual else urllib.parse.urljoin(self.base_url or "", src),
            "alt": attrs.get("alt", ""),
            "title": attrs.get("title", ""),
            "caption": caption,
            "width": width,
            "height": height,
            "score": score,
            "reasons": sorted(set(reasons)),
        }
        self.candidates.append(candidate)


def parse_int(value) -> int | None:
    try:
        return int(str(value).strip())
    except (TypeError, ValueError):
        return None


def extract_image_candidates(content: str, base_url: str | None = None, limit: int = 12) -> list[dict]:
    parser = ImageCandidateExtractor(base_url=base_url)
    parser.feed(content)
    candidates = [item for item in parser.candidates if item["score"] > 0]
    candidates.sort(key=lambda item: (-item["score"], item["src"]))
    return candidates[:limit]


def is_url(value: str) -> bool:
    return bool(URL_RE.match(value))


def is_youtube_url(value: str) -> bool:
    if not is_url(value):
        return False
    host = urllib.parse.urlparse(value).netloc.lower()
    return "youtube.com" in host or "youtu.be" in host


def youtube_video_id(value: str) -> str:
    parsed = urllib.parse.urlparse(value)
    if "youtu.be" in parsed.netloc.lower():
        return parsed.path.strip("/")
    query = urllib.parse.parse_qs(parsed.query)
    return query.get("v", [""])[0]


def normalize_extracted_text(text: str, max_chars: int = 200000) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.splitlines()]
    cleaned = "\n".join(line for line in lines if line)
    return cleaned[:max_chars]


def extracted_markdown_path(vault: Path, source_label: str, title: str, text: str) -> Path:
    digest = hashlib.sha256(f"{source_label}\n{text}".encode("utf-8", errors="ignore")).hexdigest()[:12]
    root = vault / "raw" / "sources" / "extracted"
    root.mkdir(parents=True, exist_ok=True)
    return root / f"{slugify(title or source_label)[:80]}-{digest}.md"


def write_extracted_markdown(vault: Path, source_label: str, kind: str, title: str, text: str, warnings: list[str] | None = None) -> Path:
    text = normalize_extracted_text(text)
    if not text:
        raise ValueError("NO_TEXT_EXTRACTED")
    title = title.strip() or page_title_from_path(Path(urllib.parse.urlparse(source_label).path or source_label))
    out = extracted_markdown_path(vault, source_label, title, text)
    warning_lines = [f"- {warning}" for warning in (warnings or [])] or ["- None"]
    out.write_text(
        "\n".join(
            [
                f"# {title}",
                "",
                f"Original source: `{source_label}`",
                f"Extraction kind: `{kind}`",
                f"Extracted at: `{now_iso()}`",
                "",
                "## Extraction Warnings",
                "",
                *warning_lines,
                "",
                "## Extracted Text",
                "",
                text,
                "",
            ]
        ),
        encoding="utf-8",
    )
    return out


def evidence_index_path(vault: Path) -> Path:
    return vault / "evidence" / "index.jsonl"


def timestamp_seconds(stamp: str) -> int:
    parts = [int(part) for part in stamp.split(":")]
    if len(parts) == 3:
        return parts[0] * 3600 + parts[1] * 60 + parts[2]
    if len(parts) == 2:
        return parts[0] * 60 + parts[1]
    return parts[0] if parts else 0


def make_evidence_record(
    vault: Path,
    source_label: str,
    source_path: Path,
    extraction_kind: str,
    title: str,
    evidence_kind: str,
    locator: str,
    text: str,
    location: dict | None = None,
    parser: str | None = None,
    original: dict | None = None,
) -> dict:
    clean_text = normalize_extracted_text(text, max_chars=12000)
    source_ref = vault_rel(vault, source_path)
    evidence_id = "ev_" + hashlib.sha256(f"{source_label}\n{extraction_kind}\n{locator}".encode("utf-8", errors="ignore")).hexdigest()[:20]
    record = {
        "evidence_id": evidence_id,
        "chunk_id": evidence_id,
        "evidence_kind": evidence_kind,
        "title": title,
        "source": source_label,
        "source_path": source_ref,
        "extraction_kind": extraction_kind,
        "locator": locator,
        "location": location or {},
        "text": clean_text,
        "content_sha256": sha256_text(clean_text),
        "indexed_at": now_iso(),
    }
    if is_url(source_label):
        record["source_url"] = source_label
    if parser:
        record["parser"] = parser
    if original:
        record["original"] = jsonable(original)
        if original.get("path"):
            record["original_path"] = original["path"]
        if original.get("sha256"):
            record["original_sha256"] = original["sha256"]
        if original.get("source_url"):
            record["source_url"] = original["source_url"]
    return record


def markdown_block(text: str, start: int, end_pattern: str) -> str:
    body = text[start:]
    match = re.search(end_pattern, body, flags=re.MULTILINE)
    return body[: match.start()].strip() if match else body.strip()


def extracted_evidence_records(
    vault: Path,
    source_label: str,
    source_path: Path,
    extraction_kind: str,
    title: str,
    text: str,
    parser: str | None = None,
    original: dict | None = None,
) -> list[dict]:
    records = []

    for page_match in re.finditer(r"^## Page (\d+)\s*$", text, flags=re.MULTILINE):
        page = int(page_match.group(1))
        body = markdown_block(text, page_match.end(), r"^## Page \d+\s*$")
        if body:
            records.append(make_evidence_record(vault, source_label, source_path, extraction_kind, title, "page", f"page:{page}", body, {"page": page}, parser, original))
        for table_match in re.finditer(r"^### Table (\d+)\s*$", body, flags=re.MULTILINE):
            table = int(table_match.group(1))
            table_body = markdown_block(body, table_match.end(), r"^### ")
            if table_body:
                records.append(make_evidence_record(vault, source_label, source_path, extraction_kind, title, "table", f"page:{page}:table:{table}", table_body, {"page": page, "table": table}, parser, original))

    for match in re.finditer(r"^- \[(\d{2}:\d{2}:\d{2})\]\s+(.+)$", text, flags=re.MULTILINE):
        stamp, line = match.groups()
        records.append(make_evidence_record(vault, source_label, source_path, extraction_kind, title, "timestamp", f"timestamp:{stamp}", line, {"timestamp": stamp, "seconds": timestamp_seconds(stamp)}, parser, original))

    for match in re.finditer(r"^- (\d+)\. (.*?): `([^`]+)` score=([0-9.]+) reasons=(.*)$", text, flags=re.MULTILINE):
        index, label, src, score, reasons = match.groups()
        records.append(
            make_evidence_record(
                vault,
                source_label,
                source_path,
                extraction_kind,
                title,
                "image_candidate",
                f"image-candidate:{index}",
                f"{label}: {src}",
                {"candidate_index": int(index), "src": src, "score": float(score), "reasons": [item.strip() for item in reasons.split(",") if item.strip()]},
                parser,
                original,
            )
        )

    for match in re.finditer(r"^- Frame (\d+): `([^`]+)`(.*)$", text, flags=re.MULTILINE):
        index, frame_path, rest = match.groups()
        records.append(make_evidence_record(vault, source_label, source_path, extraction_kind, title, "video_frame", f"frame:{index}", f"{frame_path} {rest}".strip(), {"frame_index": int(index), "frame_path": frame_path}, parser, original))

    if not records:
        for heading_match in re.finditer(r"^## (Slide \d+|Sheet: .+|Table \d+|Metadata|Transcript|Image Candidates|Frame Files)\s*$", text, flags=re.MULTILINE):
            heading = heading_match.group(1).strip()
            body = markdown_block(text, heading_match.end(), r"^## ")
            if body:
                records.append(make_evidence_record(vault, source_label, source_path, extraction_kind, title, "section", f"section:{slugify(heading)}", body, {"section": heading}, parser, original))

    if not records and text.strip():
        records.append(make_evidence_record(vault, source_label, source_path, extraction_kind, title, "document", "document", text, {}, parser, original))
    return records


def upsert_evidence_records(vault: Path, records: list[dict]) -> list[dict]:
    if not records:
        return []
    path = evidence_index_path(vault)
    path.parent.mkdir(parents=True, exist_ok=True)
    by_id = {}
    if path.exists():
        for line in path.read_text(encoding="utf-8", errors="ignore").splitlines():
            if not line.strip():
                continue
            try:
                item = json.loads(line)
            except json.JSONDecodeError:
                continue
            if item.get("evidence_id"):
                by_id[item["evidence_id"]] = item
    for record in records:
        by_id[record["evidence_id"]] = jsonable(record)
    path.write_text("\n".join(json.dumps(item, ensure_ascii=False, sort_keys=True) for item in by_id.values()) + "\n", encoding="utf-8")
    return records


def evidence_records_for_source_path(vault: Path, source_path: Path) -> list[dict]:
    path = evidence_index_path(vault)
    if not path.exists():
        return []
    source_ref = vault_rel(vault, source_path)
    records = []
    for line in path.read_text(encoding="utf-8", errors="ignore").splitlines():
        if not line.strip():
            continue
        try:
            item = json.loads(line)
        except json.JSONDecodeError:
            continue
        if item.get("source_path") == source_ref:
            records.append(item)
    return records


def index_extracted_evidence(
    vault: Path,
    source_label: str,
    source_path: Path,
    extraction_kind: str,
    title: str,
    text: str,
    parser: str | None = None,
    original: dict | None = None,
) -> list[dict]:
    records = extracted_evidence_records(vault, source_label, source_path, extraction_kind, title, text, parser=parser, original=original)
    return upsert_evidence_records(vault, records)


def image_average_hash(path: Path) -> str:
    try:
        from PIL import Image
    except Exception as exc:
        raise RuntimeError("Image analysis requires Pillow.") from exc
    with Image.open(path) as image:
        gray = image.convert("L").resize((8, 8))
        values = list(gray.tobytes())
    mean = sum(values) / len(values)
    bits = "".join("1" if value >= mean else "0" for value in values)
    return f"{int(bits, 2):016x}"


def image_stats(path: Path) -> dict:
    try:
        from PIL import Image, ImageStat
    except Exception as exc:
        raise RuntimeError("Image analysis requires Pillow.") from exc
    with Image.open(path) as image:
        rgb = image.convert("RGB")
        gray = image.convert("L")
        stat = ImageStat.Stat(gray)
        sample = rgb.resize((1, 1)).getpixel((0, 0))
        return {
            "width": image.width,
            "height": image.height,
            "mode": image.mode,
            "brightness": round(float(stat.mean[0]), 2),
            "contrast": round(float(stat.stddev[0]), 2),
            "average_rgb": list(sample),
            "aspect": "landscape" if image.width >= image.height else "portrait",
            "ahash": image_average_hash(path),
        }


def ocr_image(path: Path, languages: str = "eng+kor") -> tuple[str, str | None]:
    if command_available("tesseract"):
        cmd = ["tesseract", str(path), "stdout", "-l", languages, "--psm", "6"]
        result = subprocess.run(cmd, text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=45, check=False)
        if result.returncode == 0:
            return normalize_extracted_text(result.stdout, max_chars=4000), None
        return "", f"tesseract failed: {result.stderr.strip()[:300]}"
    try:
        import pytesseract
        from PIL import Image
    except Exception:
        return "", "OCR unavailable: install tesseract or pytesseract."
    try:
        with Image.open(path) as image:
            return normalize_extracted_text(pytesseract.image_to_string(image, lang=languages), max_chars=4000), None
    except Exception as exc:
        return "", f"pytesseract failed: {exc}"


def frame_asset_root(vault: Path, run_id: str) -> Path:
    root = vault / "raw" / "figures" / "video-frames" / run_id
    root.mkdir(parents=True, exist_ok=True)
    return root


def video_frame_run_id(source_label: str, title: str = "") -> str:
    basis = f"{source_label}\n{title}"
    parsed = urllib.parse.urlparse(source_label)
    if is_youtube_url(source_label):
        name = youtube_video_id(source_label) or parsed.netloc
    else:
        name = Path(parsed.path or source_label).stem or parsed.netloc or "video-frames"
    return f"{slugify(title or name)[:60]}-{hashlib.sha256(basis.encode('utf-8', errors='ignore')).hexdigest()[:10]}"


def download_youtube_video(url: str, out_dir: Path) -> tuple[Path, str, list[str]]:
    dlp = yt_dlp_command()
    if not dlp:
        raise RuntimeError("YouTube frame extraction requires yt-dlp or uvx.")
    if not command_available("ffmpeg"):
        raise RuntimeError("YouTube frame extraction requires ffmpeg.")
    output_template = str(out_dir / "source.%(ext)s")
    cmd = [
        *dlp,
        "-f",
        "bestvideo[height<=720]+bestaudio/best[height<=720]/best",
        "--merge-output-format",
        "mp4",
        "--write-info-json",
        "--print",
        "title",
        "-o",
        output_template,
        url,
    ]
    result = subprocess.run(cmd, text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=300, check=False)
    title = result.stdout.strip().splitlines()[0] if result.stdout.strip() else f"YouTube {youtube_video_id(url) or 'video'}"
    if result.returncode != 0:
        raise RuntimeError(f"yt-dlp video download failed: {result.stderr.strip()[:500]}")
    videos = sorted(path for path in out_dir.glob("source.*") if path.suffix.lower() in VIDEO_SUFFIXES)
    if not videos:
        raise RuntimeError("yt-dlp did not produce a supported video file.")
    warnings = []
    if result.stderr.strip():
        warnings.append(result.stderr.strip()[:500])
    return videos[0], title, warnings


def extract_video_frames(video_path: Path, frames_dir: Path, every_seconds: float = 30.0, max_frames: int = 12) -> list[Path]:
    if not command_available("ffmpeg"):
        raise RuntimeError("Video frame extraction requires ffmpeg.")
    if every_seconds <= 0:
        raise ValueError("every_seconds must be positive.")
    frames_dir.mkdir(parents=True, exist_ok=True)
    output = frames_dir / "frame_%06d.jpg"
    cmd = [
        "ffmpeg",
        "-hide_banner",
        "-loglevel",
        "error",
        "-y",
        "-i",
        str(video_path),
        "-vf",
        f"fps=1/{every_seconds}",
        "-frames:v",
        str(max_frames),
        str(output),
    ]
    result = subprocess.run(cmd, text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=180, check=False)
    if result.returncode != 0:
        raise RuntimeError(f"ffmpeg frame extraction failed: {result.stderr.strip()[:500]}")
    frames = sorted(frames_dir.glob("frame_*.jpg"))
    if not frames:
        raise RuntimeError("ffmpeg did not produce frames.")
    return frames


def prepare_frame_source(vault: Path, source_value: str, every_seconds: float, max_frames: int) -> tuple[str, str, Path, list[Path], list[str], dict]:
    value = str(source_value)
    warnings = []
    run_id = video_frame_run_id(value)
    root = frame_asset_root(vault, run_id)
    if is_youtube_url(value):
        video_path, title, video_warnings = download_youtube_video(value, root)
        original = preserve_original_file(vault, video_path, source_label=value)
        original["source_url"] = value
        warnings.extend(video_warnings)
        frames = extract_video_frames(video_path, root, every_seconds=every_seconds, max_frames=max_frames)
        return run_id, title, root, frames, warnings, original
    if is_url(value):
        body, content_type = fetch_url(value, max_bytes=100_000_000)
        suffix = Path(urllib.parse.urlparse(value).path).suffix.lower()
        title = page_title_from_path(Path(urllib.parse.urlparse(value).path or "remote-video"))
        original = preserve_original_bytes(vault, value, body, suffix or ".bin", source_url=value, content_type=content_type)
        if content_type.startswith("image/") or suffix in IMAGE_SUFFIXES:
            target = root / f"frame_000001{suffix if suffix in IMAGE_SUFFIXES else '.jpg'}"
            target.write_bytes(body)
            return run_id, title, root, [target], warnings, original
        if content_type.startswith("video/") or suffix in VIDEO_SUFFIXES:
            video_path = root / f"source{suffix if suffix in VIDEO_SUFFIXES else '.mp4'}"
            video_path.write_bytes(body)
            frames = extract_video_frames(video_path, root, every_seconds=every_seconds, max_frames=max_frames)
            return run_id, title, root, frames, warnings, original
        raise RuntimeError(f"Unsupported URL content type for frame analysis: {content_type or suffix}")
    path = Path(value).expanduser()
    if not path.exists() or not path.is_file():
        raise RuntimeError("Frame source is missing or not a file.")
    title = page_title_from_path(path)
    suffix = path.suffix.lower()
    original = preserve_original_file(vault, path)
    if suffix in IMAGE_SUFFIXES:
        target = root / f"frame_000001{suffix}"
        if not target.exists():
            shutil.copy2(path, target)
        return run_id, title, root, [target], warnings, original
    if suffix in VIDEO_SUFFIXES:
        frames = extract_video_frames(path, root, every_seconds=every_seconds, max_frames=max_frames)
        return run_id, title, root, frames, warnings, original
    raise RuntimeError(f"Unsupported frame source suffix: {suffix}")


def analyze_frames(vault: Path, frames: list[Path], run_id: str, ocr: bool = False, ocr_languages: str = "eng+kor") -> tuple[list[dict], list[str]]:
    results = []
    warnings = []
    for index, frame in enumerate(frames, start=1):
        rel = frame.relative_to(vault).as_posix()
        item = {"index": index, "path": rel, "bytes": frame.stat().st_size}
        try:
            item.update(image_stats(frame))
        except Exception as exc:
            warnings.append(f"{rel}: image analysis failed: {exc}")
        if ocr:
            text, warning = ocr_image(frame, languages=ocr_languages)
            item["ocr_text"] = text
            if warning:
                item["ocr_warning"] = warning
                warnings.append(f"{rel}: {warning}")
        results.append(item)
    manifest = vault / "swarmvault" / "exports" / "video-frames" / f"{run_id}.json"
    manifest.parent.mkdir(parents=True, exist_ok=True)
    manifest.write_text(json.dumps({"run_id": run_id, "frames": results, "warnings": warnings}, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return results, warnings


def write_video_frame_markdown(vault: Path, source_label: str, title: str, run_id: str, frames: list[dict], warnings: list[str], original: dict | None = None) -> Path:
    warning_lines = [f"- {warning}" for warning in warnings] or ["- None"]
    frame_lines = []
    ocr_lines = []
    for frame in frames:
        frame_lines.append(
            "- Frame {index}: `{path}` size={width}x{height} brightness={brightness} contrast={contrast} ahash={ahash}".format(
                index=frame.get("index"),
                path=frame.get("path"),
                width=frame.get("width", "?"),
                height=frame.get("height", "?"),
                brightness=frame.get("brightness", "?"),
                contrast=frame.get("contrast", "?"),
                ahash=frame.get("ahash", "?"),
            )
        )
        if frame.get("ocr_text"):
            ocr_lines.append(f"### Frame {frame.get('index')} OCR\n\n{frame['ocr_text'][:2000]}\n")
    text = "\n".join(
        [
            f"# {title} Video Frame Analysis",
            "",
            f"Original source: `{source_label}`",
            "Extraction kind: `video-frames`",
            f"Frame run: `{run_id}`",
            f"Extracted at: `{now_iso()}`",
            "",
            "## Extraction Warnings",
            "",
            *warning_lines,
            "",
            "## Frame Files",
            "",
            *frame_lines,
            "",
            "## Visual Notes",
            "",
            "- Use these frames as visual evidence for chart, figure, slide, experiment, or trading strategy review.",
            "- Compare OCR text and frame timing with transcript notes when a YouTube transcript is also ingested.",
            "- Reuse chart-like frames through `figure-card`, `figure-insight`, and `figure-export`.",
            "",
            "## OCR Text",
            "",
            *(ocr_lines or ["- OCR was not run or no readable text was detected."]),
            "",
        ]
    )
    path = write_extracted_markdown(vault, source_label, "video-frames", f"{title} Video Frames", text, warnings=warnings)
    index_extracted_evidence(vault, source_label, path, "video-frames", f"{title} Video Frames", text, original=original)
    return path


def video_frames_command(
    vault: Path,
    source: str,
    every_seconds: float = 30.0,
    max_frames: int = 12,
    ocr: bool = False,
    ocr_languages: str = "eng+kor",
    ingest: bool = False,
    figure_card: bool = False,
) -> dict:
    vault = vault.expanduser()
    run_id, title, root, frame_paths, source_warnings, original = prepare_frame_source(vault, source, every_seconds, max_frames)
    frame_records, analysis_warnings = analyze_frames(vault, frame_paths, run_id, ocr=ocr, ocr_languages=ocr_languages)
    warnings = source_warnings + analysis_warnings
    source_markdown = write_video_frame_markdown(vault, source, title, run_id, frame_records, warnings, original=original)
    report = {
        "source": source,
        "run_id": run_id,
        "title": title,
        "frames_dir": root.relative_to(vault).as_posix(),
        "frames": len(frame_records),
        "source_markdown": source_markdown.relative_to(vault).as_posix(),
        "manifest": (vault / "swarmvault" / "exports" / "video-frames" / f"{run_id}.json").relative_to(vault).as_posix(),
        "original": original,
        "evidence_index": evidence_index_path(vault).relative_to(vault).as_posix(),
        "warnings": warnings,
    }
    if ingest:
        ingest_result = ingest_sources(vault, [source_markdown.as_posix()])
        report["ingest"] = ingest_result
        if figure_card:
            pages = [item["source_note"] for item in ingest_result.get("ingested", []) if item.get("source_note")]
            report["figure_card"] = figure_card_command(vault, pages, domain="auto", all_sources=False)
    return report


def collect_markdown_outputs(root: Path) -> str:
    parts = []
    for file in sorted(root.rglob("*.md")):
        text = file.read_text(encoding="utf-8", errors="ignore").strip()
        if text:
            parts.append(f"## {file.relative_to(root).as_posix()}\n\n{text}")
    return "\n\n".join(parts)


def run_pdf_external_tool(path: Path, tool: str) -> tuple[str, list[str]]:
    warnings = []
    with tempfile.TemporaryDirectory(prefix=f"ghpf-{tool}-") as temp_dir:
        out_dir = Path(temp_dir) / "out"
        out_dir.mkdir(parents=True, exist_ok=True)
        if tool == "opendataloader-pdf":
            cmd = ["opendataloader-pdf", str(path), "-o", str(out_dir), "-f", "markdown"]
        elif tool == "marker_single":
            cmd = ["marker_single", str(path), "--output_format", "markdown", "--output_dir", str(out_dir)]
        else:
            raise ValueError(f"Unsupported PDF tool: {tool}")
        completed = subprocess.run(cmd, text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False, timeout=300)
        if completed.returncode != 0:
            raise RuntimeError(f"{tool} failed: {completed.stderr.strip()[:500]}")
        if completed.stderr.strip():
            warnings.append(f"{tool}: {completed.stderr.strip()[:300]}")
        text = collect_markdown_outputs(out_dir)
        if not text.strip():
            raise RuntimeError(f"{tool} produced no Markdown output.")
        return text, warnings


def extract_pdf_with_pdfplumber(path: Path) -> tuple[str, list[str]]:
    try:
        import pdfplumber
    except Exception as exc:
        raise RuntimeError("pdfplumber is not available.") from exc
    parts = []
    table_count = 0
    with pdfplumber.open(str(path)) as pdf:
        for page_index, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            page_parts = [f"## Page {page_index}"]
            if text.strip():
                page_parts.append(text.strip())
            for table_index, table in enumerate(page.extract_tables() or [], start=1):
                if not table:
                    continue
                table_count += 1
                rows = []
                for row in table:
                    cells = [re.sub(r"\s+", " ", str(cell or "")).strip() for cell in row]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    page_parts.append(f"### Table {table_index}\n\n" + "\n".join(rows))
            parts.append("\n\n".join(page_parts))
    warnings = [f"pdfplumber extracted {table_count} tables."] if table_count else []
    return "\n\n".join(parts), warnings


def extract_pdf_with_pypdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except Exception:
        try:
            from PyPDF2 import PdfReader
        except Exception as exc:
            raise RuntimeError("PDF extraction requires pypdf or PyPDF2.") from exc
    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        if page_text.strip():
            pages.append(f"## Page {index}\n\n{page_text.strip()}")
    return "\n\n".join(pages)


def extract_pdf_text(path: Path) -> str:
    return extract_pdf_with_pypdf(path)


def extract_pdf_document(path: Path) -> tuple[str, list[str], str]:
    attempts = []
    for tool in ("opendataloader-pdf", "marker_single"):
        if not command_available(tool):
            attempts.append(f"{tool}: unavailable")
            continue
        try:
            text, warnings = run_pdf_external_tool(path, tool)
            return text, warnings + attempts, tool
        except Exception as exc:
            attempts.append(f"{tool}: {exc}")
    if python_module_available("pdfplumber"):
        try:
            text, warnings = extract_pdf_with_pdfplumber(path)
            return text, warnings + attempts, "pdfplumber"
        except Exception as exc:
            attempts.append(f"pdfplumber: {exc}")
    text = extract_pdf_with_pypdf(path)
    return text, attempts, "pypdf"


def extract_docx_text(path: Path) -> str:
    from docx import Document
    document = Document(str(path))
    parts = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table_index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            parts.append(f"## Table {table_index}\n\n" + "\n".join(rows))
    return "\n\n".join(parts)


def extract_pptx_text(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text_frame.text.replace("\n", " ").strip() for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                slide_parts.append("\n".join(rows))
        if slide_parts:
            parts.append(f"## Slide {slide_index}\n\n" + "\n\n".join(slide_parts))
    return "\n\n".join(parts)


def extract_xlsx_text(path: Path) -> str:
    try:
        import openpyxl
    except Exception as exc:
        raise RuntimeError("XLSX extraction requires openpyxl.") from exc
    workbook = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    parts = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(max_row=80, values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                rows.append("| " + " | ".join(values) + " |")
        if rows:
            parts.append(f"## Sheet: {sheet.title}\n\n" + "\n".join(rows))
    return "\n\n".join(parts)


def convert_document_with_kordoc(path: Path) -> tuple[str, list[str]]:
    if not command_available("npx"):
        raise RuntimeError("kordoc conversion requires npx.")
    with tempfile.TemporaryDirectory(prefix="ghpf-kordoc-") as temp_dir:
        out = Path(temp_dir) / "output.md"
        cmd = ["npx", "-y", "kordoc", str(path), "-o", str(out), "--silent"]
        completed = subprocess.run(cmd, text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False, timeout=240)
        if completed.returncode != 0:
            raise RuntimeError(f"kordoc failed: {completed.stderr.strip()[:500]}")
        if not out.exists() or not out.read_text(encoding="utf-8", errors="ignore").strip():
            raise RuntimeError("kordoc produced no Markdown output.")
        warnings = [completed.stderr.strip()[:300]] if completed.stderr.strip() else []
        return out.read_text(encoding="utf-8", errors="ignore"), warnings


def convert_hwp_with_hwpjs(path: Path) -> tuple[str, list[str]]:
    with tempfile.TemporaryDirectory(prefix="ghpf-hwpjs-") as temp_dir:
        out = Path(temp_dir) / "output.md"
        if command_available("hwpjs"):
            cmd = ["hwpjs", "to-markdown", str(path), "-o", str(out), "--images-dir", str(Path(temp_dir) / "images")]
        elif command_available("npx"):
            cmd = ["npx", "-y", "@ohah/hwpjs", "to-markdown", str(path), "-o", str(out), "--images-dir", str(Path(temp_dir) / "images")]
        else:
            raise RuntimeError("HWP conversion requires hwpjs or npx.")
        completed = subprocess.run(cmd, text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False, timeout=240)
        if completed.returncode != 0:
            raise RuntimeError(f"hwpjs failed: {completed.stderr.strip()[:500]}")
        if not out.exists() or not out.read_text(encoding="utf-8", errors="ignore").strip():
            raise RuntimeError("hwpjs produced no Markdown output.")
        warnings = [completed.stderr.strip()[:300]] if completed.stderr.strip() else []
        return out.read_text(encoding="utf-8", errors="ignore"), warnings


def extract_office_document(path: Path) -> tuple[str, list[str], str]:
    suffix = path.suffix.lower()
    attempts = []
    if suffix in {".hwp", ".hwpx", ".hwpml"}:
        for name, func in (("kordoc", convert_document_with_kordoc), ("hwpjs", convert_hwp_with_hwpjs)):
            try:
                text, warnings = func(path)
                return text, warnings + attempts, name
            except Exception as exc:
                attempts.append(f"{name}: {exc}")
        raise RuntimeError("; ".join(attempts) or "No HWP converter available.")
    if suffix == ".docx":
        return extract_docx_text(path), attempts, "python-docx"
    if suffix == ".pptx":
        return extract_pptx_text(path), attempts, "python-pptx"
    if suffix in {".xlsx", ".xls"}:
        if suffix == ".xls":
            try:
                text, warnings = convert_document_with_kordoc(path)
                return text, warnings + attempts, "kordoc"
            except Exception as exc:
                attempts.append(f"kordoc: {exc}")
        return extract_xlsx_text(path), attempts, "openpyxl"
    raise RuntimeError(f"Unsupported office suffix: {suffix}")


def extract_html_text(content: str) -> tuple[str, str]:
    parser = SimpleHTMLTextExtractor()
    parser.feed(content)
    return parser.title, parser.text


def extract_html_document(content: str, base_url: str | None = None) -> tuple[str, str, list[str]]:
    title, text = extract_html_text(content)
    candidates = extract_image_candidates(content, base_url=base_url)
    warnings = []
    parts = [text]
    if candidates:
        parts.extend(["", "## Image Candidates", ""])
        for index, item in enumerate(candidates, start=1):
            label = item.get("alt") or item.get("caption") or item.get("title") or item.get("tag")
            parts.append(
                f"- {index}. {label}: `{item['src']}` score={item['score']} reasons={', '.join(item['reasons']) or 'n/a'}"
            )
    else:
        warnings.append("No content-bearing image candidates detected.")
    return title, "\n".join(part for part in parts if part is not None), warnings


def decode_response_body(body: bytes, content_type: str) -> str:
    match = re.search(r"charset=([^;\s]+)", content_type, re.IGNORECASE)
    encoding = match.group(1) if match else "utf-8"
    return body.decode(encoding, errors="ignore")


def fetch_url(url: str, max_bytes: int = 10_000_000) -> tuple[bytes, str]:
    request = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0 GHFP-LLM-Wiki/1.0"})
    with urllib.request.urlopen(request, timeout=30) as response:
        return response.read(max_bytes), response.headers.get("content-type", "")


def fetch_url_with_playwright(url: str) -> tuple[str, list[str]]:
    if not python_module_available("playwright"):
        raise RuntimeError("Python playwright is not available.")
    script = r"""
import json, sys
from playwright.sync_api import sync_playwright
url = sys.argv[1]
with sync_playwright() as p:
    browser = p.chromium.launch(headless=True)
    page = browser.new_page()
    page.goto(url, wait_until="networkidle", timeout=30000)
    result = {"title": page.title(), "html": page.content()}
    browser.close()
print(json.dumps(result, ensure_ascii=False))
"""
    result = subprocess.run(["python3", "-c", script, url], text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False, timeout=60)
    if result.returncode != 0:
        raise RuntimeError(f"playwright fallback failed: {result.stderr.strip()[:300]}")
    data = json.loads(result.stdout)
    return data.get("html", ""), [f"playwright_title: {data.get('title', '')}"]


def fetch_url_with_deepcloak(url: str) -> tuple[str, list[str]]:
    if not command_available("deepcloak"):
        raise RuntimeError("deepcloak is not available.")
    with tempfile.TemporaryDirectory(prefix="ghpf-deepcloak-") as temp_dir:
        out = Path(temp_dir) / "report.md"
        query = f"Extract the main readable content from this page with citations: {url}"
        cmd = ["deepcloak", query, "--depth", "quick", "--out", str(out)]
        result = subprocess.run(cmd, text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False, timeout=180)
        if result.returncode != 0:
            raise RuntimeError(f"deepcloak fallback failed: {result.stderr.strip()[:300]}")
        if not out.exists():
            raise RuntimeError("deepcloak did not produce a report.")
        return out.read_text(encoding="utf-8", errors="ignore"), ["deepcloak fallback report"]


def fetch_web_content(url: str) -> tuple[bytes, str, list[str]]:
    warnings = []
    try:
        body, content_type = fetch_url(url)
        return body, content_type, warnings
    except Exception as exc:
        warnings.append(f"urllib fetch failed: {exc}")
    for name, func in (("playwright", fetch_url_with_playwright), ("deepcloak", fetch_url_with_deepcloak)):
        try:
            text, fallback_warnings = func(url)
            return text.encode("utf-8"), "text/html; charset=utf-8", warnings + [f"{name} fallback used"] + fallback_warnings
        except Exception as exc:
            warnings.append(f"{name} fallback failed: {exc}")
    raise RuntimeError("; ".join(warnings))


def save_downloaded_pdf(vault: Path, url: str, body: bytes) -> Path:
    digest = hashlib.sha256(body).hexdigest()[:12]
    name = slugify(Path(urllib.parse.urlparse(url).path).stem or urllib.parse.urlparse(url).netloc or "download")
    root = vault / "raw" / "sources" / "downloads"
    root.mkdir(parents=True, exist_ok=True)
    path = root / f"{name}-{digest}.pdf"
    if not path.exists():
        path.write_bytes(body)
    return path


def clean_vtt_text(text: str) -> str:
    lines = []
    previous = None
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line or line == "WEBVTT" or line.startswith(("Kind:", "Language:", "NOTE")):
            continue
        if "-->" in line or re.fullmatch(r"\d+", line):
            continue
        line = re.sub(r"<[^>]+>", "", line)
        line = html.unescape(line).strip()
        if line and line != previous:
            lines.append(line)
            previous = line
    return "\n".join(lines)


def yt_dlp_command() -> list[str] | None:
    if command_available("yt-dlp"):
        return ["yt-dlp"]
    if command_available("uvx"):
        return ["uvx", "yt-dlp"]
    return None


def seconds_to_timestamp(value: float) -> str:
    seconds = int(value or 0)
    hours, remainder = divmod(seconds, 3600)
    minutes, seconds = divmod(remainder, 60)
    return f"{hours:02d}:{minutes:02d}:{seconds:02d}"


def youtube_metadata(url: str) -> tuple[dict, list[str]]:
    dlp = yt_dlp_command()
    if not dlp:
        return {}, ["yt-dlp unavailable for metadata."]
    cmd = [*dlp, "--dump-json", "--skip-download", url]
    completed = subprocess.run(cmd, text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False, timeout=90)
    if completed.returncode != 0:
        return {}, [f"yt-dlp metadata failed: {completed.stderr.strip()[:300]}"]
    try:
        data = json.loads(completed.stdout)
    except json.JSONDecodeError as exc:
        return {}, [f"yt-dlp metadata JSON decode failed: {exc}"]
    keys = ["id", "title", "channel", "uploader", "upload_date", "duration", "view_count", "like_count", "webpage_url", "description"]
    metadata = {key: data.get(key) for key in keys if data.get(key) not in (None, "")}
    metadata["chapters"] = [
        {"start_time": chapter.get("start_time"), "title": chapter.get("title", "")}
        for chapter in data.get("chapters", [])[:50]
        if chapter.get("title")
    ]
    return metadata, []


def format_youtube_document(title: str, snippets: list[dict], metadata: dict | None = None) -> str:
    lines = [f"# {title}", ""]
    metadata = metadata or {}
    if metadata:
        lines.extend(["## Metadata", ""])
        for key, value in metadata.items():
            if key == "chapters":
                continue
            rendered = str(value).replace("\n", " ")[:1000]
            lines.append(f"- {key}: {rendered}")
        if metadata.get("chapters"):
            lines.extend(["", "## Chapters", ""])
            for chapter in metadata["chapters"]:
                lines.append(f"- [{seconds_to_timestamp(chapter.get('start_time') or 0)}] {chapter.get('title', '')}")
        lines.append("")
    lines.extend(["## Transcript", ""])
    for item in snippets:
        text = re.sub(r"\s+", " ", str(item.get("text", ""))).strip()
        if not text:
            continue
        lines.append(f"- [{seconds_to_timestamp(float(item.get('start', 0) or 0))}] {text}")
    return "\n".join(lines)


def transcript_from_youtube_api(video_id: str) -> tuple[list[dict], str]:
    from youtube_transcript_api import YouTubeTranscriptApi

    api = YouTubeTranscriptApi()
    for languages in (["ko", "en"], ["ko"], ["en"]):
        try:
            transcript = api.fetch(video_id, languages=languages)
            snippets = [{"text": item.text, "start": item.start, "duration": item.duration} for item in transcript]
            return snippets, f"youtube_transcript_api.fetch:{','.join(languages)}"
        except Exception:
            continue
    try:
        raw = YouTubeTranscriptApi.get_transcript(video_id, languages=["ko", "en"])
        snippets = [{"text": item.get("text", ""), "start": item.get("start", 0), "duration": item.get("duration", 0)} for item in raw]
        return snippets, "youtube_transcript_api.get_transcript"
    except Exception as exc:
        raise RuntimeError(str(exc)) from exc


def transcript_from_ytdlp(url: str, temp_dir: Path) -> tuple[list[dict], str]:
    dlp = yt_dlp_command()
    if not dlp:
        raise RuntimeError("yt-dlp unavailable.")
    output_template = str(temp_dir / "%(id)s.%(ext)s")
    cmd = [
        *dlp,
        "--skip-download",
        "--write-subs",
        "--write-auto-subs",
        "--sub-langs",
        "ko.*,en.*,ko,en",
        "--sub-format",
        "json3/vtt",
        "-o",
        output_template,
        url,
    ]
    result = subprocess.run(cmd, text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=120, check=False)
    if result.returncode != 0:
        raise RuntimeError(f"yt-dlp failed: {result.stderr.strip()[:500]}")
    json_files = sorted(temp_dir.glob("*.json3"))
    if json_files:
        data = json.loads(json_files[0].read_text(encoding="utf-8", errors="ignore"))
        snippets = []
        for event in data.get("events", []):
            text = "".join(seg.get("utf8", "") for seg in event.get("segs", [])).strip()
            if text:
                snippets.append({"text": text, "start": (event.get("tStartMs") or 0) / 1000, "duration": (event.get("dDurationMs") or 0) / 1000})
        if snippets:
            return snippets, "yt-dlp:json3"
    vtt_files = sorted(temp_dir.glob("*.vtt"))
    if vtt_files:
        text = "\n".join(clean_vtt_text(path.read_text(encoding="utf-8", errors="ignore")) for path in vtt_files)
        snippets = [{"text": line, "start": 0, "duration": 0} for line in text.splitlines() if line.strip()]
        if snippets:
            return snippets, "yt-dlp:vtt"
    raise RuntimeError("yt-dlp did not produce usable subtitle files.")


def extract_youtube_transcript(url: str) -> tuple[str, str, list[str]]:
    video_id = youtube_video_id(url)
    warnings = []
    metadata, metadata_warnings = youtube_metadata(url)
    warnings.extend(metadata_warnings)
    title = metadata.get("title") or f"YouTube {video_id or 'video'}"
    if video_id:
        try:
            snippets, method = transcript_from_youtube_api(video_id)
            warnings.append(f"transcript_method: {method}")
            return title, format_youtube_document(title, snippets, metadata), warnings
        except Exception as exc:
            warnings.append(f"youtube_transcript_api unavailable or failed: {exc}")

    dlp = yt_dlp_command()
    if not dlp:
        raise RuntimeError("YouTube extraction requires youtube_transcript_api, yt-dlp, or uvx.")

    with tempfile.TemporaryDirectory(prefix="ghpf-youtube-") as temp_dir:
        snippets, method = transcript_from_ytdlp(url, Path(temp_dir))
        warnings.append(f"transcript_method: {method}")
        return title, format_youtube_document(title, snippets, metadata), warnings


def extract_source_to_markdown(vault: Path, source_value: str | Path) -> dict:
    vault = vault.expanduser()
    value = str(source_value)
    if is_youtube_url(value):
        original = external_original_record(value)
        title, text, warnings = extract_youtube_transcript(value)
        path = write_extracted_markdown(vault, value, "youtube-transcript", title, text, warnings=warnings)
        records = index_extracted_evidence(vault, value, path, "youtube-transcript", title, text, original=original)
        return {"source": value, "path": path, "kind": "youtube-transcript", "extracted": True, "original": original, "evidence_index": evidence_index_path(vault), "evidence_records": len(records), "warnings": warnings}
    if is_url(value):
        body, content_type, fetch_warnings = fetch_web_content(value)
        if "application/pdf" in content_type.lower() or urllib.parse.urlparse(value).path.lower().endswith(".pdf"):
            pdf_path = save_downloaded_pdf(vault, value, body)
            original = preserve_original_file(vault, pdf_path, source_label=value)
            original["source_url"] = value
            if content_type:
                original["content_type"] = content_type
            text, warnings, method = extract_pdf_document(pdf_path)
            path = write_extracted_markdown(vault, value, f"url-pdf:{method}", page_title_from_path(pdf_path), text, warnings=fetch_warnings + warnings)
            records = index_extracted_evidence(vault, value, path, f"url-pdf:{method}", page_title_from_path(pdf_path), text, parser=method, original=original)
            return {"source": value, "path": path, "kind": "url-pdf", "parser": method, "extracted": True, "download": pdf_path, "original": original, "evidence_index": evidence_index_path(vault), "evidence_records": len(records), "warnings": fetch_warnings + warnings}
        original = preserve_original_bytes(vault, value, body, ".html", source_url=value, content_type=content_type)
        title, text, warnings = extract_html_document(decode_response_body(body, content_type), base_url=value)
        path = write_extracted_markdown(vault, value, "web-page", title or urllib.parse.urlparse(value).netloc, text, warnings=fetch_warnings + warnings)
        records = index_extracted_evidence(vault, value, path, "web-page", title or urllib.parse.urlparse(value).netloc, text, original=original)
        return {"source": value, "path": path, "kind": "web-page", "extracted": True, "original": original, "evidence_index": evidence_index_path(vault), "evidence_records": len(records), "warnings": fetch_warnings + warnings}

    path = Path(value).expanduser()
    if not path.exists() or not path.is_file():
        return {"source": value, "error": "missing_or_not_file"}
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        original = preserve_original_file(vault, path)
        text, warnings, method = extract_pdf_document(path)
        out = write_extracted_markdown(vault, str(path), f"pdf:{method}", page_title_from_path(path), text, warnings=warnings)
        records = index_extracted_evidence(vault, str(path), out, f"pdf:{method}", page_title_from_path(path), text, parser=method, original=original)
        return {"source": value, "path": out, "kind": "pdf", "parser": method, "extracted": True, "original_path": path, "original": original, "evidence_index": evidence_index_path(vault), "evidence_records": len(records), "warnings": warnings}
    if suffix in {".html", ".htm"}:
        original = preserve_original_file(vault, path)
        title, text, warnings = extract_html_document(path.read_text(encoding="utf-8", errors="ignore"), base_url=path.resolve().as_uri())
        out = write_extracted_markdown(vault, str(path), "html-file", title or page_title_from_path(path), text, warnings=warnings)
        records = index_extracted_evidence(vault, str(path), out, "html-file", title or page_title_from_path(path), text, original=original)
        return {"source": value, "path": out, "kind": "html-file", "extracted": True, "original_path": path, "original": original, "evidence_index": evidence_index_path(vault), "evidence_records": len(records), "warnings": warnings}
    if suffix in OFFICE_SOURCE_SUFFIXES:
        original = preserve_original_file(vault, path)
        text, warnings, method = extract_office_document(path)
        out = write_extracted_markdown(vault, str(path), f"office:{method}", page_title_from_path(path), text, warnings=warnings)
        records = index_extracted_evidence(vault, str(path), out, f"office:{method}", page_title_from_path(path), text, parser=method, original=original)
        return {"source": value, "path": out, "kind": "office", "parser": method, "extracted": True, "original_path": path, "original": original, "evidence_index": evidence_index_path(vault), "evidence_records": len(records), "warnings": warnings}
    linked_evidence = evidence_records_for_source_path(vault, path)
    original = linked_evidence[0].get("original") if linked_evidence and linked_evidence[0].get("original") else preserve_original_file(vault, path)
    result = {"source": value, "path": path, "kind": "file", "extracted": False, "original_path": path, "original": original}
    if linked_evidence:
        result["evidence_index"] = evidence_index_path(vault)
        result["evidence_records"] = len(linked_evidence)
    return result


def extract_sources(vault: Path, sources: list[str]) -> dict:
    extracted = []
    skipped = []
    for source in sources:
        try:
            result = extract_source_to_markdown(vault, source)
        except Exception as exc:
            skipped.append({"source": str(source), "reason": "extract_failed", "error": str(exc)})
            continue
        if result.get("error"):
            skipped.append({"source": result["source"], "reason": result["error"]})
        else:
            record = {key: jsonable(value) for key, value in result.items() if key != "original_path"}
            extracted.append(record)
    return {"extracted": extracted, "skipped": skipped}


def ingest_sources(vault: Path, sources: list[str], move: bool = False) -> dict:
    vault = vault.expanduser()
    selected = sources if sources else source_candidates(vault)
    manifest = load_manifest(vault)
    known_hashes = {item.get("sha256") for item in manifest.get("sources", [])}
    ingested = []
    skipped = []
    extracted = []

    for source_value in selected:
        try:
            prepared = extract_source_to_markdown(vault, source_value)
        except Exception as exc:
            skipped.append({"source": str(source_value), "reason": "extract_failed", "error": str(exc)})
            continue
        if prepared.get("error"):
            skipped.append({"source": prepared["source"], "reason": prepared["error"]})
            continue
        source = prepared["path"]
        if prepared.get("extracted"):
            extracted.append({key: jsonable(value) for key, value in prepared.items() if key != "original_path"})
        if not source.exists() or not source.is_file():
            skipped.append({"source": str(source), "reason": "missing_or_not_file"})
            continue
        raw_path = copy_to_raw(vault, source)
        digest = sha256_file(raw_path)
        if digest in known_hashes:
            skipped.append({"source": str(source), "reason": "already_ingested"})
            continue

        text = read_text(raw_path)
        title = page_title_from_path(raw_path)
        source_rel = raw_path.relative_to(vault).as_posix()
        linked_evidence = evidence_records_for_source_path(vault, raw_path)
        original = prepared.get("original") or {}
        if linked_evidence and linked_evidence[0].get("original"):
            original = linked_evidence[0]["original"]
        if not linked_evidence and text.strip():
            linked_evidence = index_extracted_evidence(
                vault,
                str(prepared.get("source") or source_value),
                raw_path,
                str(prepared.get("kind") or "file"),
                title,
                text,
                parser=prepared.get("parser"),
                original=original or None,
            )
        evidence_record_count = prepared.get("evidence_records") or len(linked_evidence)
        original_ref = original.get("path") or original.get("source_url") or original.get("source") or ""
        evidence_index_rel = evidence_index_path(vault).relative_to(vault).as_posix() if evidence_record_count else ""
        source_note = vault / "wiki" / "sources" / f"{slugify(raw_path.stem)}.md"
        terms = extract_terms(text)
        summary = summarize_text(text)
        term_links = " ".join(f"[[{titleize_slug(term)}]]" for term in terms[:8])
        key_points = source_key_points(text)

        source_note.parent.mkdir(parents=True, exist_ok=True)
        frontmatter = {
            "tags": ["ghpf/source"],
            "source": source_rel,
            "created": now_iso(),
            "aliases": [title],
            "sha256": digest,
        }
        if original_ref:
            frontmatter["original"] = original_ref
        if original.get("sha256"):
            frontmatter["original_sha256"] = original["sha256"]
        if evidence_index_rel:
            frontmatter["evidence_index"] = evidence_index_rel
            frontmatter["evidence_records"] = evidence_record_count
        source_note.write_text(
            frontmatter_block(frontmatter)
            +
            "\n".join(
                [
                    f"# {title}",
                    "",
                    f"Source: `{source_rel}`",
                    f"SHA256: `{digest}`",
                    *( [f"Original: `{original_ref}`"] if original_ref else [] ),
                    *( [f"Original SHA256: `{original['sha256']}`"] if original.get("sha256") else [] ),
                    *( [f"Evidence index: `{evidence_index_rel}` ({evidence_record_count} records)"] if evidence_index_rel else [] ),
                    "",
                    "## Summary",
                    "",
                    *[f"- {line}" for line in summary],
                    "",
                    "## Key Links",
                    "",
                    term_links or "No key terms extracted.",
                    "",
                    "## Source Coverage",
                    "",
                    *[f"- [x] {line}" for line in key_points],
                    "",
                    "## Open Questions",
                    "",
                    "- What should be merged into durable concept or entity pages?",
                    "",
                ]
            ),
            encoding="utf-8",
        )

        concept_pages = []
        for term in terms[:8]:
            concept_title = titleize_slug(term)
            concept_path = vault / "wiki" / "concepts" / f"{slugify(term)}.md"
            if not concept_path.exists():
                concept_path.parent.mkdir(parents=True, exist_ok=True)
                concept_path.write_text(
                    frontmatter_block(
                        {
                            "tags": ["ghpf/concept"],
                            "source": source_rel,
                            "created": now_iso(),
                            "aliases": [concept_title],
                        }
                    )
                    + f"# {concept_title}\n\n## Sources\n\n- [[{title}]]\n\n## Notes\n\n- Created from `{source_rel}`.\n",
                    encoding="utf-8",
                )
            else:
                current = concept_path.read_text(encoding="utf-8")
                source_link = f"- [[{title}]]"
                if source_link not in current:
                    with concept_path.open("a", encoding="utf-8") as handle:
                        handle.write(f"\n{source_link}\n")
            concept_pages.append(concept_path.relative_to(vault).as_posix())

        page_rel = source_note.relative_to(vault).as_posix()
        ensure_index_entry(vault, page_rel, title)
        for concept in concept_pages:
            ensure_index_entry(vault, concept, Path(concept).stem)

        manifest.setdefault("sources", []).append(
            {
                "path": source_rel,
                "sha256": digest,
                "ingested_at": now_iso(),
                "source_note": page_rel,
                "original": jsonable(original) if original else None,
                "evidence_index": evidence_index_rel or None,
                "evidence_records": evidence_record_count,
            }
        )
        manifest.setdefault("generated_pages", []).append(page_rel)
        manifest.setdefault("generated_pages", []).extend(concept_pages)
        manifest.setdefault("operations", []).append({"type": "ingest", "time": now_iso(), "source": source_rel})
        append_log(vault, f"ingested `{source_rel}` into `{page_rel}`.")
        ingested.append({"source": source_rel, "source_note": page_rel, "concept_pages": concept_pages, "original": jsonable(original) if original else None, "evidence_index": evidence_index_rel or None, "evidence_records": evidence_record_count})

        original_path = prepared.get("original_path")
        if move and isinstance(original_path, Path) and original_path.exists() and not original_path.resolve().is_relative_to((vault / "raw").resolve()):
            original_path.unlink()

    manifest["generated_pages"] = sorted(set(manifest.get("generated_pages", [])))
    save_manifest(vault, manifest)
    return {"ingested": ingested, "skipped": skipped, "extracted": extracted}


def build_graph(vault: Path) -> dict:
    nodes = {}
    edges = []
    title_to_id = {}
    files = markdown_files(vault)

    for path in files:
        rel = path.relative_to(vault).as_posix()
        title = path.stem
        nodes[rel] = {"id": rel, "title": title, "path": rel, "kind": "page"}
        for alias in page_aliases(path):
            title_to_id[alias] = rel

    for path in files:
        source = path.relative_to(vault).as_posix()
        text = path.read_text(encoding="utf-8", errors="ignore")
        for match in WIKILINK_RE.findall(text):
            target_title = match.strip()
            target = title_to_id.get(target_title.lower(), f"wiki/unresolved/{slugify(target_title)}.md")
            if target not in nodes:
                nodes[target] = {"id": target, "title": target_title, "path": target, "kind": "unresolved"}
            edges.append({"source": source, "target": target, "type": "wikilink"})

    graph = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "node_count": len(nodes),
        "edge_count": len(edges),
        "nodes": list(nodes.values()),
        "edges": edges,
    }

    state = vault / "swarmvault" / "state"
    exports = vault / "swarmvault" / "exports"
    state.mkdir(parents=True, exist_ok=True)
    exports.mkdir(parents=True, exist_ok=True)
    (state / "graph.json").write_text(json.dumps(graph, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    (exports / "graph.html").write_text(render_graph_html(graph), encoding="utf-8")
    return graph


def render_graph_html(graph: dict) -> str:
    data = html.escape(json.dumps(graph, ensure_ascii=False))
    return f"""<!doctype html>
<html lang="en">
<head><meta charset="utf-8"><title>GHFP Wiki Graph</title>
<style>body{{font-family:system-ui;margin:2rem;line-height:1.5}}pre{{white-space:pre-wrap;background:#f6f8fa;padding:1rem}}</style></head>
<body>
<h1>GHFP Wiki Graph</h1>
<p>Nodes: {graph['node_count']} | Edges: {graph['edge_count']}</p>
<pre id="graph">{data}</pre>
</body></html>
"""


def build_context(vault: Path, query: str, limit: int = 8, max_chars: int = 1200) -> Path:
    search = hybrid_search(vault, query, limit=max(limit * 4, limit, 1))
    selected = []
    seen_paths = set()
    for item in search["results"]:
        path = vault / item["path"]
        rel = item["path"]
        if path.exists() and is_retrieval_page(path) and rel not in seen_paths:
            selected.append((item["score"], path, path.read_text(encoding="utf-8", errors="ignore")))
            seen_paths.add(rel)
        if len(selected) >= limit:
            break

    out_dir = vault / "swarmvault" / "context-packs"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"{datetime.now(timezone.utc).strftime('%Y%m%dT%H%M%SZ')}-{slugify(query)[:80]}.md"

    lines = [
        f"# Context Pack: {query}",
        "",
        f"Generated: {datetime.now(timezone.utc).isoformat()}",
        "",
        "## Included Pages",
        "",
    ]
    for score, path, text in selected:
        rel = path.relative_to(vault).as_posix()
        excerpt = text.strip().replace("\n\n", "\n")[:max_chars]
        lines.extend([f"### {rel}", "", f"Score: {score}", "", "```md", excerpt, "```", ""])
    if not selected:
        lines.extend(["No matching wiki pages found.", ""])

    out_path.write_text("\n".join(lines), encoding="utf-8")
    return out_path


def record_task(vault: Path, status: str, title: str, note: str = "", target: str = "") -> Path:
    now = datetime.now(timezone.utc).isoformat()
    record = {"time": now, "status": status, "title": title, "target": target, "note": note}
    ledger = vault / "swarmvault" / "task-ledger"
    ledger.mkdir(parents=True, exist_ok=True)
    with (ledger / "tasks.jsonl").open("a", encoding="utf-8") as handle:
        handle.write(json.dumps(record, ensure_ascii=False) + "\n")

    task_dir = vault / "wiki" / "tasks"
    task_dir.mkdir(parents=True, exist_ok=True)
    task_path = task_dir / f"{slugify(title)}.md"
    if not task_path.exists():
        task_path.write_text(f"# {title}\n\n", encoding="utf-8")
    with task_path.open("a", encoding="utf-8") as handle:
        handle.write(f"- {now} `{status}`")
        if target:
            handle.write(f" target={target}")
        if note:
            handle.write(f": {note}")
        handle.write("\n")
    return task_path


def file_back(vault: Path, title: str, body: str, folder: str = "wiki/syntheses") -> Path:
    vault = vault.expanduser()
    folder_path = vault / folder
    folder_path.mkdir(parents=True, exist_ok=True)
    page_path = folder_path / f"{slugify(title)}.md"
    if page_path.exists():
        with page_path.open("a", encoding="utf-8") as handle:
            handle.write(f"\n## Update {now_iso()}\n\n{body.strip()}\n")
    else:
        page_path.write_text(
            frontmatter_block(
                {
                    "tags": ["ghpf/synthesis"],
                    "source": "file-back",
                    "created": now_iso(),
                    "aliases": [title],
                }
            )
            + f"# {title}\n\n{body.strip()}\n",
            encoding="utf-8",
        )

    rel = page_path.relative_to(vault).as_posix()
    ensure_index_entry(vault, rel, title)
    append_log(vault, f"file-backed synthesis `{rel}`.")
    manifest = load_manifest(vault)
    manifest.setdefault("generated_pages", []).append(rel)
    manifest["generated_pages"] = sorted(set(manifest.get("generated_pages", [])))
    manifest.setdefault("operations", []).append({"type": "file_back", "time": now_iso(), "page": rel})
    save_manifest(vault, manifest)
    return page_path


def quality_score(page: Path, source_text: str | None = None, consistency: float | None = None, suggestions: float | None = None, coverage: float | None = None) -> dict:
    content = page.read_text(encoding="utf-8", errors="ignore")
    frontmatter, body = parse_frontmatter(content)
    filled = sum(1 for field in QUALITY_REQUIRED_FIELDS if frontmatter.get(field))
    completeness = filled / len(QUALITY_REQUIRED_FIELDS)
    link_count = len(WIKILINK_RE.findall(body))
    if link_count == 0:
        connections = 0.2
    elif link_count <= 2:
        connections = 0.6
    elif link_count <= 8:
        connections = 1.0
    else:
        connections = 0.8

    coverage_value = coverage
    if coverage_value is None:
        checked = len(re.findall(r"- \[[xX]\]", body))
        unchecked = len(re.findall(r"- \[ \]", body))
        if checked or unchecked:
            coverage_value = checked / max(checked + unchecked, 1)
        elif source_text:
            source_terms = set(extract_terms(source_text, limit=20))
            body_terms = set(extract_terms(body, limit=80))
            coverage_value = len(source_terms & body_terms) / max(len(source_terms), 1)
        else:
            coverage_value = 0.5

    consistency_value = 0.5 if consistency is None else consistency
    suggestions_value = 0.5 if suggestions is None else suggestions
    breakdown = {
        "completeness": {"value": round(completeness, 3), "weight": QUALITY_WEIGHTS["completeness"], "filled": filled, "required": QUALITY_REQUIRED_FIELDS},
        "connections": {"value": round(connections, 3), "weight": QUALITY_WEIGHTS["connections"], "wikilinks": link_count, "target_range": "3-8"},
        "coverage": {"value": round(coverage_value, 3), "weight": QUALITY_WEIGHTS["coverage"]},
        "consistency": {"value": round(consistency_value, 3), "weight": QUALITY_WEIGHTS["consistency"]},
        "suggestions": {"value": round(suggestions_value, 3), "weight": QUALITY_WEIGHTS["suggestions"]},
    }
    score = round(sum(item["value"] * item["weight"] for item in breakdown.values()), 3)
    return {"page": str(page), "lint_score": score, "passed": score >= 0.7, "threshold": 0.7, "breakdown": breakdown}


def run_quality(vault: Path, pages: list[str], strict: bool = False, consistency: float | None = None, suggestions: float | None = None, coverage: float | None = None) -> dict:
    vault = vault.expanduser()
    selected = [vault / page for page in pages] if pages else [p for p in markdown_files(vault) if is_retrieval_page(p)]
    results = []
    for page in selected:
        if page.exists() and page.is_file():
            results.append(quality_score(page, consistency=consistency, suggestions=suggestions, coverage=coverage))
        else:
            results.append({"page": str(page), "error": "FILE_NOT_FOUND", "passed": False})
    report = {
        "ok": all(item.get("passed") for item in results) if strict else True,
        "strict": strict,
        "results": results,
        "average_score": round(sum(item.get("lint_score", 0.0) for item in results) / max(len(results), 1), 3),
    }
    out = vault / "swarmvault" / "exports" / "quality-report.json"
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return report


def vault_relative_path(vault: Path, value: str) -> Path:
    path = Path(value).expanduser()
    if path.is_absolute():
        return path
    return vault.expanduser() / path


def extract_sections(content: str) -> dict[str, str]:
    _, body = parse_frontmatter(content)
    sections = {}
    current = "(intro)"
    lines = []
    for line in body.splitlines():
        if re.match(r"^##\s+", line):
            if lines:
                sections[current] = "\n".join(lines).strip()
            current = line.strip()
            lines = []
        else:
            lines.append(line)
    if lines:
        sections[current] = "\n".join(lines).strip()
    return sections


def section_diff(existing: Path, new: Path) -> dict:
    old_sections = extract_sections(existing.read_text(encoding="utf-8", errors="ignore"))
    new_sections = extract_sections(new.read_text(encoding="utf-8", errors="ignore"))
    results = []
    matched_new = set()

    for old_heading, old_body in old_sections.items():
        if old_heading == "(intro)":
            continue
        if old_heading in new_sections:
            matched_new.add(old_heading)
            similarity = SequenceMatcher(None, old_body, new_sections[old_heading]).ratio()
            if similarity < 1.0:
                results.append({"type": "CHANGED", "heading": old_heading, "similarity": round(similarity, 3)})
            continue
        best_heading = None
        best_score = 0.0
        for heading, body in new_sections.items():
            if heading in matched_new or heading == "(intro)":
                continue
            heading_score = SequenceMatcher(None, old_heading, heading).ratio()
            body_score = SequenceMatcher(None, old_body, body).ratio()
            score = max(heading_score, body_score * 0.8)
            if score >= 0.4 and body_score >= 0.5 and score > best_score:
                best_heading = heading
                best_score = score
        if best_heading:
            matched_new.add(best_heading)
            results.append({"type": "RENAMED", "old_heading": old_heading, "heading": best_heading, "similarity": round(best_score, 3)})
        else:
            results.append({"type": "REMOVED", "heading": old_heading})

    for heading, body in new_sections.items():
        if heading not in matched_new and heading != "(intro)":
            results.append({"type": "NEW", "heading": heading, "preview": body[:160].replace("\n", " ")})

    return {
        "summary": {
            "new": sum(1 for item in results if item["type"] == "NEW"),
            "changed": sum(1 for item in results if item["type"] == "CHANGED"),
            "removed": sum(1 for item in results if item["type"] == "REMOVED"),
            "renamed": sum(1 for item in results if item["type"] == "RENAMED"),
        },
        "sections": results,
    }


def state_path(vault: Path) -> Path:
    return vault / "swarmvault" / "state" / "pipeline-state.json"


def load_state(vault: Path) -> dict:
    path = state_path(vault)
    if not path.exists():
        return {"created_at": now_iso(), "completed": {}, "current": PIPELINE_STEPS[0], "errors": []}
    return json.loads(path.read_text(encoding="utf-8"))


def save_state(vault: Path, state: dict) -> None:
    path = state_path(vault)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(state, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def state_action(vault: Path, action: str, step: str | None = None) -> dict:
    vault = vault.expanduser()
    if action == "init":
        state = {"created_at": now_iso(), "completed": {}, "current": PIPELINE_STEPS[0], "errors": []}
        save_state(vault, state)
        return state
    state = load_state(vault)
    if action == "show":
        return {**state, "progress": f"{len(state.get('completed', {}))}/{len(PIPELINE_STEPS)}"}
    if step not in PIPELINE_STEPS:
        return {"error": "INVALID_STEP", "valid_steps": PIPELINE_STEPS}
    step_index = PIPELINE_STEPS.index(step)
    missing = [prior for prior in PIPELINE_STEPS[:step_index] if prior not in state.get("completed", {})]
    if action == "check":
        return {"step": step, "ready": not missing, "missing": missing, "completed": list(state.get("completed", {}).keys())}
    if action == "complete":
        if missing:
            error = {"error": "STEP_SKIP_DETECTED", "step": step, "missing": missing}
            state.setdefault("errors", []).append({**error, "time": now_iso()})
            save_state(vault, state)
            return error
        state.setdefault("completed", {})[step] = now_iso()
        state["current"] = PIPELINE_STEPS[step_index + 1] if step_index + 1 < len(PIPELINE_STEPS) else "done"
        save_state(vault, state)
        return state
    return {"error": "INVALID_ACTION"}


def graph_index(vault: Path) -> tuple[dict[str, Path], dict[str, set[str]], dict[str, set[str]]]:
    pages = markdown_files(vault)
    page_by_title = {}
    for page in pages:
        for alias in page_aliases(page):
            page_by_title[alias] = page
    outgoing = {p.relative_to(vault).as_posix(): set() for p in pages}
    incoming = {p.relative_to(vault).as_posix(): set() for p in pages}
    for page in pages:
        source_rel = page.relative_to(vault).as_posix()
        text = page.read_text(encoding="utf-8", errors="ignore")
        for target in WIKILINK_RE.findall(text):
            target_page = page_by_title.get(target.strip().lower())
            if target_page:
                target_rel = target_page.relative_to(vault).as_posix()
                outgoing[source_rel].add(target_rel)
                incoming.setdefault(target_rel, set()).add(source_rel)
    return page_by_title, outgoing, incoming


def link_audit(vault: Path) -> dict:
    lint = lint_wiki(vault)
    _, outgoing, incoming = graph_index(vault.expanduser())
    total_links = sum(len(v) for v in outgoing.values())
    one_way = []
    for source, targets in outgoing.items():
        for target in targets:
            if source not in outgoing.get(target, set()):
                one_way.append({"source": source, "target": target})
    report = {
        "pages": len(outgoing),
        "links": total_links,
        "average_links_per_page": round(total_links / max(len(outgoing), 1), 3),
        "broken_links": lint["broken_links"],
        "orphans": [page for page in outgoing if not outgoing.get(page) and not incoming.get(page)],
        "deadends": [page for page in outgoing if not outgoing.get(page) and incoming.get(page)],
        "hubs": sorted([page for page, links in outgoing.items() if len(links) + len(incoming.get(page, set())) >= 10]),
        "one_way_links": one_way,
    }
    out = vault.expanduser() / "swarmvault" / "exports" / "link-audit.json"
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return report


def page_terms(page: Path) -> set[str]:
    text = page.read_text(encoding="utf-8", errors="ignore")
    return set(extract_terms(text, limit=60))


def link_strengthen(vault: Path, page_rel: str, max_links: int = 5, backlink: bool = False) -> dict:
    vault = vault.expanduser()
    page = vault / page_rel
    if not page.exists():
        return {"error": "FILE_NOT_FOUND", "page": page_rel}
    target_terms = page_terms(page)
    existing_text = page.read_text(encoding="utf-8", errors="ignore")
    existing_links = {link.lower() for link in WIKILINK_RE.findall(existing_text)}
    candidates = []
    for other in markdown_files(vault):
        if other == page or not is_retrieval_page(other):
            continue
        overlap = target_terms & page_terms(other)
        if not overlap or other.stem.lower() in existing_links or titleize_slug(other.stem).lower() in existing_links:
            continue
        candidates.append((len(overlap), other, sorted(overlap)[:8]))
    candidates.sort(key=lambda item: (-item[0], item[1].as_posix()))
    selected = candidates[:max_links]
    if selected:
        with page.open("a", encoding="utf-8") as handle:
            handle.write("\n## Related Notes\n\n")
            for score, other, overlap in selected:
                title = titleize_slug(other.stem)
                handle.write(f"- [[{title}]] - overlap: {', '.join(overlap)}\n")
        if backlink:
            source_title = titleize_slug(page.stem)
            for _, other, _ in selected:
                text = other.read_text(encoding="utf-8", errors="ignore")
                link = f"[[{source_title}]]"
                if link not in text:
                    with other.open("a", encoding="utf-8") as handle:
                        handle.write(f"\n## Backlinks\n\n- {link}\n")
    report = {
        "page": page_rel,
        "links_added": [{"page": other.relative_to(vault).as_posix(), "score": score, "overlap": overlap} for score, other, overlap in selected],
        "backlink": backlink,
    }
    out = vault / "swarmvault" / "exports" / "link-strengthen.json"
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return report


def command_available(command: str) -> bool:
    return shutil.which(command) is not None


def python_module_available(module: str) -> bool:
    result = subprocess.run(
        ["python3", "-c", f"import {module}"],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        check=False,
    )
    return result.returncode == 0


def capabilities(vault: Path | None = None) -> dict:
    modules = ["pypdf", "PyPDF2", "pdfplumber", "docx", "pptx", "openpyxl", "networkx", "youtube_transcript_api", "pytesseract", "playwright", "matplotlib", "numpy", "PIL"]
    caps = {
        "commands": {name: command_available(name) for name in ["git", "node", "npx", "uv", "uvx", "graphify", "playwright", "yt-dlp", "ffmpeg", "ffprobe", "tesseract", "obsidian", "opendataloader-pdf", "marker_single", "hwpjs", "deepcloak"]},
        "python_modules": {name: python_module_available(name) for name in modules},
        "optional_modes": {
            "basic_wikilink_search": True,
            "graph_sidecar": True,
            "graphify_import_ready": True,
            "graphify_cli_ready": command_available("graphify") or command_available("uv"),
            "youtube_ingest_ready": command_available("yt-dlp") or command_available("uvx") or python_module_available("youtube_transcript_api"),
            "ocr_ready": command_available("tesseract") or python_module_available("pytesseract"),
            "advanced_pdf_extract_ready": command_available("opendataloader-pdf") or command_available("marker_single") or python_module_available("pdfplumber"),
            "hwp_extract_ready": command_available("hwpjs") or command_available("npx"),
            "office_extract_ready": python_module_available("docx") or python_module_available("pptx") or python_module_available("openpyxl") or command_available("npx"),
            "playwright_ready": command_available("playwright") or python_module_available("playwright"),
            "deepcloak_ready": command_available("deepcloak"),
            "web_fallback_ready": python_module_available("playwright") or command_available("deepcloak"),
            "figure_export_ready": python_module_available("matplotlib") and python_module_available("numpy"),
            "image_frame_analysis_ready": python_module_available("PIL"),
            "video_frame_extract_ready": command_available("ffmpeg"),
            "youtube_frame_extract_ready": command_available("ffmpeg") and (command_available("yt-dlp") or command_available("uvx")),
        },
    }
    if vault is not None:
        caps["vault"] = doctor(vault)
    return caps


def edge_endpoint(edge, names: tuple[str, ...], index: int) -> str | None:
    if isinstance(edge, dict):
        for name in names:
            value = edge.get(name)
            if isinstance(value, dict):
                value = value.get("id") or value.get("name") or value.get("label")
            if value is not None:
                return str(value)
    elif isinstance(edge, (list, tuple)) and len(edge) > index:
        return str(edge[index])
    return None


def normalize_graphify_graph(data: dict) -> tuple[list[dict], list[dict]]:
    if isinstance(data.get("graph"), dict) and not data.get("nodes"):
        data = data["graph"]
    raw_nodes = data.get("nodes") or data.get("vertices") or []
    raw_edges = data.get("edges") or data.get("links") or []
    if isinstance(raw_nodes, dict):
        raw_nodes = [{**value, "id": key} if isinstance(value, dict) else {"id": key, "label": value} for key, value in raw_nodes.items()]

    nodes = {}
    for item in raw_nodes:
        if isinstance(item, dict):
            node_id = str(item.get("id") or item.get("key") or item.get("name") or item.get("label") or item.get("title") or len(nodes))
            title = str(item.get("title") or item.get("label") or item.get("name") or node_id)
            body = item.get("summary") or item.get("description") or item.get("text") or item.get("content") or ""
            kind = item.get("type") or item.get("kind") or item.get("group") or "graph-node"
            nodes[node_id] = {"id": node_id, "title": title, "body": str(body), "kind": str(kind), "raw": item}
        else:
            node_id = str(item)
            nodes[node_id] = {"id": node_id, "title": node_id, "body": "", "kind": "graph-node", "raw": {"id": node_id}}

    edges = []
    for item in raw_edges:
        source = edge_endpoint(item, ("source", "from", "start", "src"), 0)
        target = edge_endpoint(item, ("target", "to", "end", "dst"), 1)
        if not source or not target:
            continue
        relation = "related"
        if isinstance(item, dict):
            relation = str(item.get("label") or item.get("type") or item.get("relation") or "related")
        edges.append({"source": source, "target": target, "relation": relation})
        for node_id in (source, target):
            nodes.setdefault(node_id, {"id": node_id, "title": node_id, "body": "", "kind": "graph-node", "raw": {"id": node_id}})
    return list(nodes.values()), edges


def graphify_import(vault: Path, graph_json: Path, run_id: str | None = None, report: Path | None = None, html_file: Path | None = None, max_links: int = 20) -> dict:
    vault = vault.expanduser()
    graph_json = graph_json.expanduser()
    data = json.loads(graph_json.read_text(encoding="utf-8"))
    nodes, edges = normalize_graphify_graph(data)
    run_id = slugify(run_id or f"{datetime.now(timezone.utc).strftime('%Y%m%dT%H%M%SZ')}-{graph_json.parent.name or graph_json.stem}")
    import_root = vault / "graph_imports" / run_id
    node_root = import_root / "nodes"
    node_root.mkdir(parents=True, exist_ok=True)

    used_slugs = set()
    id_to_title = {}
    id_to_slug = {}
    for node in nodes:
        base = slugify(node["title"])
        slug = base
        suffix = 2
        while slug in used_slugs:
            slug = f"{base}-{suffix}"
            suffix += 1
        used_slugs.add(slug)
        id_to_title[node["id"]] = titleize_slug(slug)
        id_to_slug[node["id"]] = slug

    adjacency = {node["id"]: [] for node in nodes}
    for edge in edges:
        adjacency.setdefault(edge["source"], []).append((edge["target"], edge["relation"]))
        adjacency.setdefault(edge["target"], []).append((edge["source"], edge["relation"]))

    pages = []
    for node in nodes:
        title = id_to_title[node["id"]]
        page = node_root / f"{id_to_slug[node['id']]}.md"
        related = []
        for target, relation in adjacency.get(node["id"], [])[:max_links]:
            if target in id_to_title:
                related.append(f"- [[{id_to_title[target]}]] - {relation}")
        raw_preview = json.dumps(node["raw"], ensure_ascii=False, indent=2)[:2000]
        page.write_text(
            frontmatter_block(
                {
                    "tags": ["ghpf/graph-import"],
                    "source": graph_json.as_posix(),
                    "created": now_iso(),
                    "aliases": [title],
                    "graphify_run": run_id,
                }
            )
            + "\n".join(
                [
                    f"# {title}",
                    "",
                    "Graphify reference note. Promote durable knowledge into `wiki/` before treating it as canonical.",
                    "",
                    f"- Source node id: `{node['id']}`",
                    f"- Node kind: `{node['kind']}`",
                    "",
                    "## Summary",
                    "",
                    node["body"] or "No summary provided by graph source.",
                    "",
                    "## Related Graph Nodes",
                    "",
                    *(related or ["No related nodes imported."]),
                    "",
                    "## Raw Node Preview",
                    "",
                    "```json",
                    raw_preview,
                    "```",
                    "",
                ]
            ),
            encoding="utf-8",
        )
        pages.append(page.relative_to(vault).as_posix())

    index = import_root / "index.md"
    index.write_text(
        "\n".join(
            [
                f"# Graphify Import: {run_id}",
                "",
                "This folder is a reference layer for broad map context. Keep canonical notes in `wiki/`.",
                "",
                f"- Source graph: `{graph_json}`",
                f"- Imported nodes: {len(nodes)}",
                f"- Imported edges: {len(edges)}",
                "",
                "## Nodes",
                "",
                *[f"- [[{id_to_title[node['id']]}]]" for node in nodes],
                "",
            ]
        ),
        encoding="utf-8",
    )

    cache_root = vault / "swarmvault" / "cache" / "graphify" / run_id
    export_root = vault / "swarmvault" / "exports" / "graphify" / run_id
    cache_root.mkdir(parents=True, exist_ok=True)
    export_root.mkdir(parents=True, exist_ok=True)
    shutil.copy2(graph_json, cache_root / "graph.json")
    for optional, target_name in ((report, "GRAPH_REPORT.md"), (html_file, "graph.html")):
        if optional and optional.expanduser().exists():
            shutil.copy2(optional.expanduser(), export_root / target_name)

    manifest = load_manifest(vault)
    manifest.setdefault("graph_imports", []).append({"run_id": run_id, "graph": str(graph_json), "imported_at": now_iso(), "nodes": len(nodes), "edges": len(edges), "path": import_root.relative_to(vault).as_posix()})
    manifest.setdefault("operations", []).append({"type": "graphify_import", "time": now_iso(), "run_id": run_id})
    save_manifest(vault, manifest)
    append_log(vault, f"imported Graphify graph `{graph_json}` into `{import_root.relative_to(vault).as_posix()}`.")
    return {"run_id": run_id, "nodes": len(nodes), "edges": len(edges), "import_root": import_root.relative_to(vault).as_posix(), "pages": pages, "cache": cache_root.relative_to(vault).as_posix()}


def cache_clean(vault: Path, max_age_days: int = 30, keep_latest: int = 10, dry_run: bool = False, include_exports: bool = False) -> dict:
    vault = vault.expanduser()
    cutoff = datetime.now(timezone.utc) - timedelta(days=max_age_days)
    roots = [vault / "swarmvault" / "cache"]
    if include_exports:
        roots.append(vault / "swarmvault" / "exports" / "graphify")
    removed = []
    kept = []

    for root in roots:
        if not root.exists():
            continue
        candidates = []
        for child in root.iterdir():
            if child.is_file():
                candidates.append(child)
            elif child.is_dir():
                candidates.extend(sorted(child.iterdir()))
        candidates = [item for item in candidates if item.exists()]
        candidates.sort(key=lambda item: item.stat().st_mtime, reverse=True)
        protected = set(candidates[:keep_latest])
        for item in candidates:
            modified = datetime.fromtimestamp(item.stat().st_mtime, tz=timezone.utc)
            rel = item.relative_to(vault).as_posix()
            if item in protected or modified >= cutoff:
                kept.append(rel)
                continue
            removed.append(rel)
            if not dry_run:
                if item.is_dir():
                    shutil.rmtree(item)
                else:
                    item.unlink()
    return {"dry_run": dry_run, "max_age_days": max_age_days, "keep_latest": keep_latest, "include_exports": include_exports, "removed": removed, "kept": kept}


def maintenance_state_path(vault: Path) -> Path:
    return vault / "swarmvault" / "state" / "maintenance-state.json"


def load_maintenance_state(vault: Path) -> dict:
    path = maintenance_state_path(vault)
    if not path.exists():
        return {"created_at": now_iso(), "graphify": {}}
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return {"created_at": now_iso(), "graphify": {}, "warnings": ["maintenance-state.json was invalid JSON and was reset"]}


def save_maintenance_state(vault: Path, state: dict) -> None:
    path = maintenance_state_path(vault)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(state, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def count_sources_for_maintenance(vault: Path, manifest: dict) -> int:
    sources = manifest.get("sources", [])
    if sources:
        return len(sources)
    source_dir = vault / "wiki" / "sources"
    return len(list(source_dir.glob("*.md"))) if source_dir.exists() else 0


def summarize_lint(report: dict) -> dict:
    return {
        "ok": report.get("ok"),
        "pages": report.get("pages"),
        "broken_links_count": len(report.get("broken_links") or []),
        "missing_required": report.get("missing_required") or [],
        "manifest_missing_count": len(report.get("manifest_missing") or []),
    }


def summarize_link_audit(report: dict) -> dict:
    return {
        "pages": report.get("pages"),
        "links": report.get("links"),
        "average_links_per_page": report.get("average_links_per_page"),
        "broken_links_count": len(report.get("broken_links") or []),
        "orphans_count": len(report.get("orphans") or []),
        "deadends_count": len(report.get("deadends") or []),
        "one_way_links_count": len(report.get("one_way_links") or []),
    }


def graphify_next_steps(vault: Path, graphify_needed: bool, graphify_status: str, threshold: int) -> dict:
    graphify_input_dir = vault / "raw" / "graphify_articles"
    graphify_export_dir = vault / "swarmvault" / "exports" / "graphify"
    input_files = sorted(p for p in graphify_input_dir.rglob("*") if p.is_file()) if graphify_input_dir.exists() else []
    if graphify_status == "imported":
        return {
            "needed": False,
            "summary": "Graphify graph.json was imported and the source-count checkpoint was updated.",
            "input_dir": graphify_input_dir.relative_to(vault).as_posix(),
            "input_files": len(input_files),
            "next": [],
        }
    if not graphify_needed:
        return {
            "needed": False,
            "summary": f"Graphify is below threshold; run maintenance again after {threshold} new source notes or use --force-graphify.",
            "input_dir": graphify_input_dir.relative_to(vault).as_posix(),
            "input_files": len(input_files),
            "next": [],
        }
    return {
        "needed": True,
        "summary": "Graphify is recommended, but GHFP does not guess an external Graphify CLI command. Generate graph.json with your chosen Graphify tool, then import it.",
        "input_dir": graphify_input_dir.relative_to(vault).as_posix(),
        "input_files": len(input_files),
        "fallback_corpus": "wiki/sources",
        "expected_output_dir": graphify_export_dir.relative_to(vault).as_posix(),
        "next": [
            f"Place or generate the Graphify corpus under `{graphify_input_dir.relative_to(vault).as_posix()}` when using a bulk Graphify workflow.",
            "Run your chosen Graphify tool externally so it writes a graph.json file.",
            "Import the result with `python3 scripts/ghpf_wiki.py maintenance --vault <vault> --auto-graphify --graphify-graph <graph.json>`.",
        ],
    }


def maintenance(
    vault: Path,
    threshold: int = 20,
    auto_graphify: bool = False,
    force_graphify: bool = False,
    graphify_graph: Path | None = None,
    graphify_run_id: str | None = None,
) -> dict:
    vault = vault.expanduser()
    manifest = load_manifest(vault)
    state = load_maintenance_state(vault)
    graphify_state = state.setdefault("graphify", {})

    source_count = count_sources_for_maintenance(vault, manifest)
    last_graphify_source_count = int(graphify_state.get("last_source_count") or 0)
    new_sources_since_graphify = max(source_count - last_graphify_source_count, 0)
    graphify_needed = force_graphify or (source_count > 0 and new_sources_since_graphify >= threshold)

    index_report = build_hybrid_index(vault)
    link_report = link_audit(vault)
    graph_report = build_graph(vault)
    lint_report = lint_wiki(vault)

    graphify_status = "not_needed"
    graphify_result = None
    graphify_warning = None
    if graphify_needed and auto_graphify:
        if graphify_graph is not None:
            graphify_result = graphify_import(vault, graphify_graph, run_id=graphify_run_id)
            index_report = build_hybrid_index(vault)
            graphify_status = "imported"
            graphify_state["last_source_count"] = source_count
            graphify_state["last_run_at"] = now_iso()
            graphify_state["last_run_id"] = graphify_result.get("run_id")
            graphify_state["last_import_root"] = graphify_result.get("import_root")
            new_sources_since_graphify = 0
            graphify_needed = False
        elif command_available("graphify"):
            graphify_status = "needed_but_no_graph_json"
            graphify_warning = "graphify command is available, but this maintenance command imports graph.json and does not guess external CLI arguments"
        else:
            graphify_status = "needed_but_no_graph_json"
            graphify_warning = "provide --graphify-graph <graph.json> or run Graphify externally before import"
    elif graphify_needed:
        graphify_status = "recommended"

    state["last_checked_at"] = now_iso()
    state["source_count"] = source_count
    state["threshold"] = threshold
    state["new_sources_since_graphify"] = new_sources_since_graphify
    state["graphify_needed"] = graphify_needed
    state["last_maintenance"] = {
        "index_chunks": index_report.get("chunks"),
        "graph_nodes": graph_report.get("node_count"),
        "graph_edges": graph_report.get("edge_count"),
        "lint_ok": lint_report.get("ok"),
        "graphify_status": graphify_status,
    }
    save_maintenance_state(vault, state)

    return {
        "vault": str(vault),
        "state_path": maintenance_state_path(vault).relative_to(vault).as_posix(),
        "source_count": source_count,
        "threshold": threshold,
        "last_graphify_source_count": graphify_state.get("last_source_count", last_graphify_source_count),
        "new_sources_since_graphify": new_sources_since_graphify,
        "graphify_needed": graphify_needed,
        "graphify_status": graphify_status,
        "graphify_warning": graphify_warning,
        "graphify_next_steps": graphify_next_steps(vault, graphify_needed, graphify_status, threshold),
        "graphify_result": graphify_result,
        "index": {"chunks": index_report.get("chunks"), "backend": index_report.get("backend")},
        "link_audit": summarize_link_audit(link_report),
        "graph": {"node_count": graph_report.get("node_count"), "edge_count": graph_report.get("edge_count")},
        "lint": summarize_lint(lint_report),
    }


def page_heading(content: str, fallback: str) -> str:
    _, body = parse_frontmatter(content)
    for line in body.splitlines():
        if line.startswith("# "):
            return line[2:].strip()
    return titleize_slug(Path(fallback).stem)


def sentence_candidates(text: str, keywords: list[str], limit: int = 4) -> list[str]:
    _, body = parse_frontmatter(text)
    pieces = re.split(r"(?<=[.!?。])\s+|\n+", body)
    scored = []
    for piece in pieces:
        cleaned = re.sub(r"\s+", " ", piece).strip(" -*#\t")
        if len(cleaned) < 20:
            continue
        haystack = cleaned.lower()
        score = sum(haystack.count(keyword.lower()) for keyword in keywords)
        if score:
            scored.append((score, cleaned[:320]))
    scored.sort(key=lambda item: (-item[0], item[1]))
    return [item for _, item in scored[:limit]]


def fallback_bullets(text: str, limit: int = 4) -> list[str]:
    return [line[:320] for line in summarize_text(text, max_lines=limit)]


def card_path(vault: Path, kind: str, title: str) -> Path:
    folder = {"paper": "papers", "experiment": "experiments", "strategy": "strategies"}[kind]
    root = vault / "wiki" / "cards" / folder
    root.mkdir(parents=True, exist_ok=True)
    return root / f"{slugify(title)[:90]}.md"


def card_sections(kind: str, content: str) -> dict[str, list[str]]:
    if kind == "paper":
        return {
            "Problem": sentence_candidates(content, ["problem", "challenge", "issue", "goal", "objective", "문제", "목표"]) or fallback_bullets(content, 2),
            "Method": sentence_candidates(content, ["method", "approach", "agent", "tool", "memory", "framework", "방법", "구조"]),
            "Contribution": sentence_candidates(content, ["contribution", "improve", "shows", "proposes", "emphasizes", "기여", "제안"]),
            "Useful For My Work": sentence_candidates(content, ["irrigation", "greenhouse", "trading", "strategy", "experiment", "automatic", "관수", "자동매매"]),
            "Experiment Ideas": sentence_candidates(content, ["compare", "evaluation", "metric", "baseline", "ablation", "실험", "평가"]),
            "Limitations": sentence_candidates(content, ["limitation", "risk", "failure", "requires", "safety", "한계", "위험"]),
        }
    if kind == "experiment":
        return {
            "Hypothesis": sentence_candidates(content, ["hypothesis", "goal", "should", "improve", "검증", "가설"]) or fallback_bullets(content, 2),
            "Variables": sentence_candidates(content, ["variable", "factor", "sensor", "weather", "memory", "tool", "regime", "변수"]),
            "Metrics": sentence_candidates(content, ["metric", "accuracy", "error", "drawdown", "water", "drain", "EC", "평가", "지표"]),
            "Baselines": sentence_candidates(content, ["baseline", "control", "compare", "rule", "buy-and-hold", "기준"]),
            "Risks": sentence_candidates(content, ["risk", "failure", "leakage", "overfitting", "safety", "위험"]),
        }
    return {
        "Idea": sentence_candidates(content, ["strategy", "signal", "regime", "momentum", "volatility", "idea", "전략"]) or fallback_bullets(content, 2),
        "Components": sentence_candidates(content, ["filter", "sizing", "entry", "exit", "risk", "component", "signal"]),
        "Backtest Plan": sentence_candidates(content, ["backtest", "walk-forward", "baseline", "cost", "slippage", "validation"]),
        "Rejection Rules": sentence_candidates(content, ["reject", "drawdown", "overfit", "unstable", "leakage", "risk"]),
        "Risks": sentence_candidates(content, ["risk", "leakage", "overfitting", "transaction", "drawdown", "위험"]),
    }


def generate_card(vault: Path, page_rel: str, kind: str) -> dict:
    vault = vault.expanduser()
    page = vault_relative_path(vault, page_rel)
    if not page.exists():
        return {"page": page_rel, "error": "FILE_NOT_FOUND"}
    content = page.read_text(encoding="utf-8", errors="ignore")
    title = page_heading(content, page.name)
    terms = extract_terms(content, limit=10)
    links, candidates = existing_wikilinks_for_terms(vault, terms, limit=8)
    sections = card_sections(kind, content)
    out = card_path(vault, kind, title)
    lines = [
        frontmatter_block({"tags": [f"ghpf/card/{kind}"], "source": page.relative_to(vault).as_posix(), "created": now_iso(), "aliases": [title]}).rstrip(),
        f"# {kind.title()} Card: {title}",
        "",
        f"Source: `{page.relative_to(vault).as_posix()}`",
        "",
    ]
    for heading, bullets in sections.items():
        lines.extend([f"## {heading}", ""])
        lines.extend([f"- {bullet}" for bullet in bullets] or ["- Not enough evidence extracted yet."])
        lines.append("")
    lines.extend(["## Links", "", " ".join(links) or "No existing wiki targets found.", ""])
    if candidates:
        lines.extend(["## Candidate Terms", "", ", ".join(candidates), ""])
    out.write_text("\n".join(lines), encoding="utf-8")
    ensure_index_entry(vault, out.relative_to(vault).as_posix(), f"{kind.title()} Card: {title}")
    return {"page": page.relative_to(vault).as_posix(), "card": out.relative_to(vault).as_posix(), "kind": kind}


def card_command(vault: Path, kind: str, pages: list[str], all_sources: bool = False) -> dict:
    vault = vault.expanduser()
    if all_sources:
        selected = [p.relative_to(vault).as_posix() for p in markdown_files(vault) if p.as_posix().find("/wiki/sources/") >= 0]
    else:
        selected = pages
    results = [generate_card(vault, page, kind) for page in selected]
    out = vault / "swarmvault" / "exports" / f"{kind}-card-report.json"
    out.parent.mkdir(parents=True, exist_ok=True)
    report = {"kind": kind, "count": len(results), "results": results}
    out.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return report


def tokenize(text: str) -> list[str]:
    return [token.lower() for token in re.findall(r"[A-Za-z][A-Za-z0-9_-]{2,}|[가-힣]{2,}", text)]


def hashed_vector(text: str, dims: int = VECTOR_DIMS) -> list[float]:
    vector = [0.0] * dims
    for token in tokenize(text):
        bucket = int(hashlib.sha256(token.encode("utf-8")).hexdigest()[:8], 16) % dims
        vector[bucket] += 1.0
    norm = sum(value * value for value in vector) ** 0.5
    if norm:
        vector = [value / norm for value in vector]
    return vector


def cosine(a: list[float], b: list[float]) -> float:
    return sum(x * y for x, y in zip(a, b))


def chunk_page(vault: Path, page: Path, max_chars: int = 1200) -> list[dict]:
    content = page.read_text(encoding="utf-8", errors="ignore")
    title = page_heading(content, page.name)
    sections = extract_sections(content)
    chunks = []
    for section, body in sections.items():
        if section in NON_RETRIEVAL_SECTIONS:
            continue
        paragraphs = [p.strip() for p in re.split(r"\n\s*\n", body) if p.strip()]
        current = ""
        for paragraph in paragraphs:
            if len(current) + len(paragraph) > max_chars and current:
                chunks.append({"path": page.relative_to(vault).as_posix(), "title": title, "section": section, "text": current.strip()})
                current = ""
            current += ("\n\n" if current else "") + paragraph
        if current.strip():
            chunks.append({"path": page.relative_to(vault).as_posix(), "title": title, "section": section, "text": current.strip()})
    return chunks


def index_paths(vault: Path) -> tuple[Path, Path]:
    root = vault / "swarmvault" / "state" / "hybrid-index"
    root.mkdir(parents=True, exist_ok=True)
    return root / "chunks.jsonl", root / "vectors.sqlite"


def build_hybrid_index(vault: Path) -> dict:
    vault = vault.expanduser()
    chunks_path, db_path = index_paths(vault)
    chunks = []
    for page in reference_markdown_files(vault):
        if not is_retrieval_page(page):
            continue
        chunks.extend(chunk_page(vault, page))
    with chunks_path.open("w", encoding="utf-8") as handle:
        for idx, chunk in enumerate(chunks):
            chunk["chunk_id"] = f"chunk-{idx:06d}"
            chunk["terms"] = sorted(set(tokenize(f"{chunk['title']} {chunk['section']} {chunk['text']}")))[:80]
            handle.write(json.dumps(chunk, ensure_ascii=False) + "\n")
    with sqlite3.connect(db_path) as conn:
        conn.execute("DROP TABLE IF EXISTS vectors")
        conn.execute("CREATE TABLE vectors (chunk_id TEXT PRIMARY KEY, path TEXT, title TEXT, section TEXT, text TEXT, vector TEXT)")
        for chunk in chunks:
            conn.execute(
                "INSERT INTO vectors VALUES (?, ?, ?, ?, ?, ?)",
                (chunk["chunk_id"], chunk["path"], chunk["title"], chunk["section"], chunk["text"], json.dumps(hashed_vector(f"{chunk['title']} {chunk['section']} {chunk['text']}"))),
            )
        conn.commit()
    return {"chunks": len(chunks), "chunks_path": chunks_path.relative_to(vault).as_posix(), "vectors_path": db_path.relative_to(vault).as_posix(), "backend": "local-hashed-vector"}


def load_index_chunks(vault: Path) -> list[dict]:
    chunks_path, _ = index_paths(vault)
    if not chunks_path.exists():
        build_hybrid_index(vault)
    return [json.loads(line) for line in chunks_path.read_text(encoding="utf-8").splitlines() if line.strip()]


def hybrid_search(vault: Path, query: str, limit: int = 10) -> dict:
    vault = vault.expanduser()
    chunks = load_index_chunks(vault)
    query_terms = tokenize(query)
    query_vec = hashed_vector(query)
    scored = []
    for chunk in chunks:
        haystack = f"{chunk['title']} {chunk['section']} {chunk['text']}".lower()
        keyword = sum(haystack.count(term) for term in query_terms)
        vector = cosine(query_vec, hashed_vector(haystack))
        card_boost = 0.35 if "/cards/" in chunk["path"] else 0.0
        graph_boost = min(len(WIKILINK_RE.findall(chunk["text"])) * 0.03, 0.3)
        score = keyword + vector * 3.0 + card_boost + graph_boost
        if score > 0:
            scored.append({"score": round(score, 4), "keyword": keyword, "vector": round(vector, 4), "graph": round(graph_boost, 4), "card_boost": card_boost, **chunk})
    scored.sort(key=lambda item: (-item["score"], item["path"], item["chunk_id"]))
    results = scored[:limit]
    out = vault / "swarmvault" / "exports" / "search-report.json"
    out.parent.mkdir(parents=True, exist_ok=True)
    report = {"query": query, "limit": limit, "results": results}
    out.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return report


def insight_sections(kind: str) -> list[str]:
    if kind == "paper":
        return ["Evidence Base", "Research Gap", "Design Insight", "Related Work Angle", "Next Paper Moves"]
    if kind == "experiment":
        return ["Hypothesis", "Variables", "Baselines", "Metrics", "Ablations", "Failure Risks"]
    return ["Strategy Thesis", "Components", "Backtest Plan", "Risk Controls", "Rejection Rules"]


def evaluate_evidence(kind: str, results: list[dict]) -> dict:
    paths = {item["path"] for item in results}
    text = "\n".join(item["text"].lower() for item in results)
    checks = {
        "evidence_count": len(results),
        "source_diversity": len(paths),
        "has_method_or_component": any(word in text for word in ["method", "component", "tool", "agent", "signal", "방법", "구조"]),
        "has_metric_or_validation": any(word in text for word in ["metric", "baseline", "evaluation", "backtest", "validation", "평가", "검증"]),
        "has_risk_or_limitation": any(word in text for word in ["risk", "limitation", "drawdown", "safety", "failure", "leakage", "위험", "한계"]),
    }
    if kind == "strategy":
        checks["has_backtestability"] = any(word in text for word in ["backtest", "walk-forward", "fee", "slippage", "drawdown"])
    if kind == "experiment":
        checks["has_ablation_hint"] = any(word in text for word in ["ablation", "compare", "baseline", "variable"])
    passed = sum(1 for value in checks.values() if bool(value))
    return {"checks": checks, "score": round(passed / len(checks), 3), "passed": passed >= max(3, len(checks) - 2)}


def insight_command(vault: Path, kind: str, query: str, limit: int = 8) -> dict:
    vault = vault.expanduser()
    search = hybrid_search(vault, query, limit=limit)
    results = search["results"]
    evaluation = evaluate_evidence(kind, results)
    title = f"{kind.title()} Insight - {query[:80]}"
    out = vault / "wiki" / "syntheses" / f"{slugify(title)[:100]}.md"
    out.parent.mkdir(parents=True, exist_ok=True)
    lines = [
        frontmatter_block({"tags": [f"ghpf/insight/{kind}"], "source": "hybrid-search", "created": now_iso(), "aliases": [title]}).rstrip(),
        f"# {title}",
        "",
        f"Query: `{query}`",
        f"Evidence score: `{evaluation['score']}`",
        "",
    ]
    for heading in insight_sections(kind):
        lines.extend([f"## {heading}", ""])
        if results:
            for item in results[: min(4, len(results))]:
                snippet = re.sub(r"\s+", " ", item["text"]).strip()[:260]
                lines.append(f"- `{item['path']}`: {snippet}")
        else:
            lines.append("- No evidence found.")
        lines.append("")
    lines.extend(["## Evidence", ""])
    for item in results:
        lines.append(f"- score={item['score']} vector={item['vector']} keyword={item['keyword']} `{item['path']}` section={item['section']}")
    lines.extend(["", "## Evaluation", "", "```json", json.dumps(evaluation, ensure_ascii=False, indent=2), "```", ""])
    out.write_text("\n".join(lines), encoding="utf-8")
    ensure_index_entry(vault, out.relative_to(vault).as_posix(), title)
    return {"insight": out.relative_to(vault).as_posix(), "kind": kind, "query": query, "evidence_count": len(results), "evaluation": evaluation}


def evaluate_command(vault: Path, kind: str, target: str) -> dict:
    vault = vault.expanduser()
    path = vault_relative_path(vault, target)
    if not path.exists():
        return {"target": target, "error": "FILE_NOT_FOUND"}
    content = path.read_text(encoding="utf-8", errors="ignore")
    result = evaluate_evidence(kind, [{"path": path.relative_to(vault).as_posix(), "text": content}])
    out = vault / "swarmvault" / "evaluations" / f"{slugify(path.stem)}-{kind}.json"
    out.parent.mkdir(parents=True, exist_ok=True)
    report = {"target": path.relative_to(vault).as_posix(), "kind": kind, **result}
    out.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return report


def infer_figure_domain(text: str, fallback: str = "generic") -> str:
    lowered = text.lower()
    if any(word in lowered for word in ["btc", "bitcoin", "equity curve", "drawdown", "regime", "momentum", "volatility", "backtest", "strategy"]):
        return "trading"
    if any(word in lowered for word in ["irrigation", "greenhouse", "substrate", "drain", "ec", "water content", "sensor", "관수"]):
        return "irrigation"
    return fallback


def figure_panels(domain: str) -> list[dict]:
    if domain == "trading":
        return [
            {"label": "A", "title": "Price And Regime", "plot": "line + regime shading", "purpose": "Show when the strategy is allowed to act."},
            {"label": "B", "title": "Signal Stack", "plot": "momentum and volatility filter lines", "purpose": "Separate entry signal from risk filter."},
            {"label": "C", "title": "Position Exposure", "plot": "step plot or filled exposure", "purpose": "Show whether the strategy actually changes risk."},
            {"label": "D", "title": "Equity Curve", "plot": "strategy vs benchmark line", "purpose": "Compare absolute performance."},
            {"label": "E", "title": "Drawdown", "plot": "area plot", "purpose": "Expose failure periods and rejection evidence."},
        ]
    if domain == "irrigation":
        return [
            {"label": "A", "title": "Substrate Water Content", "plot": "time-series line", "purpose": "Show plant-water state stability."},
            {"label": "B", "title": "Irrigation Events", "plot": "bar/event markers", "purpose": "Show agent action timing and volume."},
            {"label": "C", "title": "Drainage And EC Response", "plot": "dual line or paired lines", "purpose": "Connect water action to salinity/drain response."},
            {"label": "D", "title": "Agent Decision Trace", "plot": "categorical step plot", "purpose": "Make LLM/tool decisions auditable."},
            {"label": "E", "title": "Safety Interventions", "plot": "event markers", "purpose": "Show guardrail activation and rule violations."},
        ]
    return [
        {"label": "A", "title": "Primary Outcome", "plot": "line or bar", "purpose": "Show the main dependent variable."},
        {"label": "B", "title": "Mechanism Or Signal", "plot": "line/scatter", "purpose": "Show why the outcome changed."},
        {"label": "C", "title": "Comparison", "plot": "treatment vs baseline", "purpose": "Support the claim against a baseline."},
        {"label": "D", "title": "Failure Or Uncertainty", "plot": "error/drawdown/residual", "purpose": "Expose limits and robustness."},
        {"label": "E", "title": "Decision Or Annotation Layer", "plot": "event markers", "purpose": "Make key events auditable."},
    ]


def figure_card_path(vault: Path, title: str) -> Path:
    root = vault / "wiki" / "cards" / "figures"
    root.mkdir(parents=True, exist_ok=True)
    return root / f"{slugify(title)[:90]}.md"


def generate_figure_card(vault: Path, page_rel: str, domain: str = "auto") -> dict:
    vault = vault.expanduser()
    page = vault_relative_path(vault, page_rel)
    if not page.exists():
        return {"page": page_rel, "error": "FILE_NOT_FOUND"}
    content = page.read_text(encoding="utf-8", errors="ignore")
    source_title = page_heading(content, page.name)
    resolved_domain = infer_figure_domain(content, fallback="generic") if domain == "auto" else domain
    terms = extract_terms(content, limit=12)
    links, candidates = existing_wikilinks_for_terms(vault, terms, limit=8)
    panels = figure_panels(resolved_domain)
    title = f"{source_title} Figure Pattern"
    out = figure_card_path(vault, title)
    evidence = sentence_candidates(content, ["figure", "chart", "metric", "baseline", "drawdown", "irrigation", "signal", "decision", "experiment", "backtest"], limit=5) or fallback_bullets(content, 3)
    lines = [
        frontmatter_block({"tags": ["ghpf/card/figure", f"ghpf/figure/{resolved_domain}"], "source": page.relative_to(vault).as_posix(), "created": now_iso(), "aliases": [title]}).rstrip(),
        f"# Figure Card: {title}",
        "",
        f"Source: `{page.relative_to(vault).as_posix()}`",
        f"Domain: `{resolved_domain}`",
        "",
        "## Figure Type",
        "",
        f"- {panels[0]['plot']} multi-panel diagnostic figure",
        "",
        "## Encodes",
        "",
        *[f"- Panel {panel['label']}: {panel['title']} - {panel['purpose']}" for panel in panels],
        "",
        "## Design Pattern",
        "",
        "- Use shared x-axis when panels are time-based.",
        "- Keep labels outside dense data regions.",
        "- Export vector PDF/SVG first and PNG only as a preview.",
        "",
        "## Evidence Notes",
        "",
        *[f"- {item}" for item in evidence],
        "",
        "## Useful For",
        "",
        "- Manuscript result figure",
        "- Experiment diagnostic figure",
        "- Strategy validation chart",
        "",
        "## Links",
        "",
        " ".join(links) or "No existing wiki targets found.",
        "",
    ]
    if candidates:
        lines.extend(["## Candidate Terms", "", ", ".join(candidates), ""])
    out.write_text("\n".join(lines), encoding="utf-8")
    ensure_index_entry(vault, out.relative_to(vault).as_posix(), f"Figure Card: {title}")
    return {"page": page.relative_to(vault).as_posix(), "figure_card": out.relative_to(vault).as_posix(), "domain": resolved_domain}


def figure_card_command(vault: Path, pages: list[str], domain: str = "auto", all_sources: bool = False) -> dict:
    vault = vault.expanduser()
    selected = [p.relative_to(vault).as_posix() for p in markdown_files(vault) if "/wiki/sources/" in p.as_posix()] if all_sources else pages
    results = [generate_figure_card(vault, page, domain=domain) for page in selected]
    report = {"domain": domain, "count": len(results), "results": results}
    out = vault / "swarmvault" / "exports" / "figure-card-report.json"
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return report


def read_data_schema(path_value: str | None) -> dict:
    if not path_value:
        return {}
    path = Path(path_value).expanduser()
    if not path.exists():
        return {"path": str(path), "error": "FILE_NOT_FOUND"}
    if path.suffix.lower() == ".json":
        return json.loads(path.read_text(encoding="utf-8"))
    if path.suffix.lower() == ".csv":
        first = path.read_text(encoding="utf-8", errors="ignore").splitlines()[0] if path.read_text(encoding="utf-8", errors="ignore").splitlines() else ""
        return {"path": str(path), "columns": [item.strip() for item in first.split(",") if item.strip()]}
    return {"path": str(path), "description": path.read_text(encoding="utf-8", errors="ignore")[:2000]}


def figure_insight_command(vault: Path, query: str, domain: str = "auto", data_schema: str | None = None, limit: int = 8) -> dict:
    vault = vault.expanduser()
    resolved_domain = infer_figure_domain(query, fallback="generic") if domain == "auto" else domain
    schema = read_data_schema(data_schema)
    search = hybrid_search(vault, f"{query} figure chart panel design {resolved_domain}", limit=limit)
    panels = figure_panels(resolved_domain)
    title = f"Figure Design - {query[:80]}"
    out = vault / "wiki" / "figure-designs" / f"{slugify(title)[:100]}.md"
    out.parent.mkdir(parents=True, exist_ok=True)
    lines = [
        frontmatter_block({"tags": ["ghpf/figure-design", f"ghpf/figure/{resolved_domain}"], "source": "figure-insight", "created": now_iso(), "aliases": [title]}).rstrip(),
        f"# {title}",
        "",
        f"Query: `{query}`",
        f"Domain: `{resolved_domain}`",
        "",
        "## Recommended Figure",
        "",
        f"- {len(panels)}-panel {resolved_domain} diagnostic figure",
        "- Final-size first: use double-column width for dense multi-panel figures.",
        "- Export as PDF and SVG; create PNG preview at 600 dpi.",
        "",
        "## Data Schema",
        "",
        "```json",
        json.dumps(schema, ensure_ascii=False, indent=2),
        "```",
        "",
        "## Panel Plan",
        "",
    ]
    for panel in panels:
        lines.extend([f"### Panel {panel['label']}: {panel['title']}", "", f"- Plot: {panel['plot']}", f"- Purpose: {panel['purpose']}", ""])
    lines.extend(["## Evidence From Wiki", ""])
    for item in search["results"]:
        snippet = re.sub(r"\s+", " ", item["text"]).strip()[:240]
        lines.append(f"- score={item['score']} `{item['path']}`: {snippet}")
    lines.extend(
        [
            "",
            "## Export Guidance",
            "",
            "- Keep caption text in the manuscript, not inside the image.",
            "- Use restrained colors with one accent color for intervention or regime markers.",
            "- Use panel labels A-E and shared x-axis where possible.",
            "",
        ]
    )
    out.write_text("\n".join(lines), encoding="utf-8")
    ensure_index_entry(vault, out.relative_to(vault).as_posix(), title)
    return {"figure_design": out.relative_to(vault).as_posix(), "domain": resolved_domain, "evidence_count": len(search["results"]), "panels": len(panels)}


def matplotlib_script(domain: str, figure_name: str) -> str:
    panel_count = len(figure_panels(domain))
    return f'''#!/usr/bin/env python3
from pathlib import Path
import argparse
import csv
import math

import numpy as np
import matplotlib.pyplot as plt


def mm_to_in(value):
    return value / 25.4


def read_csv(path):
    if not path:
        return {{}}
    with open(path, newline="", encoding="utf-8") as handle:
        rows = list(csv.DictReader(handle))
    data = {{}}
    for key in rows[0].keys() if rows else []:
        values = []
        for row in rows:
            try:
                values.append(float(row[key]))
            except (TypeError, ValueError):
                values.append(np.nan)
        data[key] = np.array(values, dtype=float)
    return data


def series(data, names, fallback):
    for name in names:
        if name in data and len(data[name]):
            return data[name]
    return fallback


def make_data(domain, csv_path=None):
    data = read_csv(csv_path)
    n = max((len(v) for v in data.values()), default=160)
    x = series(data, ["time", "day", "date", "index"], np.arange(n))
    t = np.arange(len(x))
    if domain == "trading":
        price = series(data, ["price", "close", "BTC", "btc"], 100 + np.cumsum(np.sin(t / 8) + np.random.default_rng(7).normal(0, 0.4, len(t))))
        momentum = series(data, ["momentum", "signal"], np.tanh(np.gradient(price) * 4))
        volatility = series(data, ["volatility", "vol"], np.clip(np.abs(np.gradient(price)) / max(np.nanstd(price), 1e-6), 0, 2))
        position = series(data, ["position", "exposure"], (momentum > 0).astype(float) * (volatility < 1.3))
        equity = series(data, ["equity", "strategy_equity"], 1 + np.cumsum(np.nan_to_num(position[:-1], nan=0, posinf=0, neginf=0).tolist() + [0]) * np.nan_to_num(np.gradient(price), nan=0) / max(np.nanmean(price), 1e-6))
        drawdown = equity / np.maximum.accumulate(equity) - 1
        return x, [price, momentum, volatility, position, equity, drawdown]
    swc = series(data, ["swc", "substrate_water_content", "water_content"], 55 + 4 * np.sin(t / 18) - 0.015 * t)
    irrigation = series(data, ["irrigation", "irrigation_volume", "volume"], np.where((t % 24) == 8, 1.2 + 0.3 * np.sin(t / 13), 0))
    drain = series(data, ["drain", "drain_fraction"], np.clip(irrigation * 0.25 + 0.05 * np.sin(t / 11), 0, None))
    ec = series(data, ["ec", "EC"], 2.8 + 0.2 * np.sin(t / 17) - 0.1 * drain)
    decision = series(data, ["decision", "agent_decision"], np.where(irrigation > 0, 1, 0))
    safety = series(data, ["safety", "guardrail", "violation"], np.where((swc < np.nanpercentile(swc, 15)) | (ec > 3.0), 1, 0))
    return x, [swc, irrigation, drain, ec, decision, safety]


def build_figure(domain, csv_path=None):
    plt.rcParams.update({{
        "font.family": "DejaVu Sans",
        "font.size": 8,
        "axes.labelsize": 8,
        "axes.titlesize": 9,
        "xtick.labelsize": 7,
        "ytick.labelsize": 7,
        "legend.fontsize": 7,
        "lines.linewidth": 1.2,
    }})
    x, values = make_data(domain, csv_path)
    height_mm = 165 if domain in ("trading", "irrigation") else 125
    fig, axes = plt.subplots({panel_count}, 1, figsize=(mm_to_in(180), mm_to_in(height_mm)), sharex=True, constrained_layout=True)
    axes = np.atleast_1d(axes)
    accent = "#2f6fbb"
    danger = "#b23a48"
    neutral = "#4c4c4c"
    if domain == "trading":
        price, momentum, volatility, position, equity, drawdown = values
        regime = momentum > 0
        axes[0].plot(x, price, color=neutral, label="BTC price")
        axes[0].fill_between(x, np.nanmin(price), np.nanmax(price), where=regime, color=accent, alpha=0.10, label="risk-on")
        axes[1].plot(x, momentum, color=accent, label="momentum")
        axes[1].plot(x, volatility, color=danger, label="volatility")
        axes[2].step(x, position, where="post", color=neutral, label="position")
        axes[3].plot(x, equity, color=accent, label="strategy equity")
        axes[4].fill_between(x, drawdown, 0, color=danger, alpha=0.35, label="drawdown")
        titles = ["A. Price and regime", "B. Signal stack", "C. Position exposure", "D. Equity curve", "E. Drawdown"]
        ylabels = ["Price", "Signal", "Exposure", "Equity", "Drawdown"]
    else:
        swc, irrigation, drain, ec, decision, safety = values
        axes[0].plot(x, swc, color=accent, label="substrate water content")
        axes[1].bar(x, irrigation, color=neutral, alpha=0.55, label="irrigation volume")
        axes[2].plot(x, drain, color=accent, label="drain fraction")
        axes[2].plot(x, ec, color=danger, label="EC")
        axes[3].step(x, decision, where="post", color=neutral, label="agent decision")
        axes[4].scatter(x[safety > 0], safety[safety > 0], color=danger, s=12, label="safety marker")
        titles = ["A. Substrate water content", "B. Irrigation events", "C. Drainage and EC response", "D. Agent decision trace", "E. Safety interventions"]
        ylabels = ["SWC (%)", "Volume", "Drain/EC", "Decision", "Event"]
    for ax, title, ylabel in zip(axes, titles, ylabels):
        ax.set_title(title, loc="left", fontweight="bold")
        ax.set_ylabel(ylabel)
        ax.grid(True, color="#d9d9d9", linewidth=0.5, alpha=0.7)
        ax.legend(frameon=False, loc="best")
    axes[-1].set_xlabel("Time")
    return fig


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--domain", default="{domain}", choices=["irrigation", "trading", "generic"])
    parser.add_argument("--csv")
    parser.add_argument("--out-dir", default=".")
    parser.add_argument("--name", default="{figure_name}")
    args = parser.parse_args()
    out_dir = Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    fig = build_figure(args.domain, args.csv)
    for ext in ("pdf", "svg", "png"):
        dpi = 600 if ext == "png" else None
        fig.savefig(out_dir / f"{{args.name}}.{{ext}}", dpi=dpi, bbox_inches="tight")
    print(out_dir / f"{{args.name}}.pdf")


if __name__ == "__main__":
    main()
'''


def figure_export_command(vault: Path, design: str | None, domain: str = "auto", name: str = "Figure_1", data: str | None = None, run: bool = False) -> dict:
    vault = vault.expanduser()
    design_text = ""
    if design:
        design_path = vault_relative_path(vault, design)
        if design_path.exists():
            design_text = design_path.read_text(encoding="utf-8", errors="ignore")
    resolved_domain = infer_figure_domain(f"{domain} {design_text}", fallback="generic") if domain == "auto" else domain
    figure_name = slugify(name).replace("-", "_") or "Figure_1"
    out_dir = vault / "swarmvault" / "exports" / "figures" / figure_name
    out_dir.mkdir(parents=True, exist_ok=True)
    script_path = out_dir / "figure.py"
    script_path.write_text(matplotlib_script(resolved_domain, figure_name), encoding="utf-8")
    result = {"domain": resolved_domain, "script": script_path.relative_to(vault).as_posix(), "out_dir": out_dir.relative_to(vault).as_posix(), "ran": False, "outputs": []}
    if run:
        cmd = ["python3", str(script_path), "--domain", resolved_domain, "--out-dir", str(out_dir), "--name", figure_name]
        if data:
            cmd.extend(["--csv", str(Path(data).expanduser())])
        completed = subprocess.run(cmd, text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False, timeout=60)
        result["ran"] = True
        result["returncode"] = completed.returncode
        result["stdout"] = completed.stdout.strip()
        result["stderr"] = completed.stderr.strip()[:1000]
        result["outputs"] = [p.relative_to(vault).as_posix() for p in sorted(out_dir.glob(f"{figure_name}.*"))]
        if completed.returncode != 0:
            result["error"] = "FIGURE_SCRIPT_FAILED"
    return result


def lint_wiki(vault: Path) -> dict:
    vault = vault.expanduser()
    pages = markdown_files(vault)
    page_by_title = {}
    for page in pages:
        for alias in page_aliases(page):
            page_by_title[alias] = page
    page_rels = {p.relative_to(vault).as_posix() for p in pages}
    broken_links = []
    orphan_pages = []
    index_missing = []
    manifest_missing = []

    linked_targets = set()
    for page in pages:
        text = page.read_text(encoding="utf-8", errors="ignore")
        for target in WIKILINK_RE.findall(text):
            title = target.strip()
            linked_targets.add(title.lower())
            if title.lower() not in page_by_title:
                broken_links.append({"page": page.relative_to(vault).as_posix(), "target": title})

    for page in pages:
        rel = page.relative_to(vault).as_posix()
        if not is_retrieval_page(page):
            continue
        if page.stem.lower() not in linked_targets and rel != "wiki/index.md":
            orphan_pages.append(rel)

    index_path = vault / "wiki" / "index.md"
    index_text = index_path.read_text(encoding="utf-8", errors="ignore") if index_path.exists() else ""
    for rel in sorted(page_rels):
        if rel.startswith("wiki/") and rel not in index_text and is_retrieval_page(Path(rel)):
            index_missing.append(rel)

    manifest = load_manifest(vault)
    for rel in manifest.get("generated_pages", []):
        if rel not in page_rels:
            manifest_missing.append(rel)

    required = ["raw", "wiki", "schema", "schema/AGENTS.md", "wiki/index.md", "wiki/log.md", "wiki/manifest.json"]
    missing_required = [item for item in required if not (vault / item).exists()]
    report = {
        "ok": not (broken_links or manifest_missing or missing_required),
        "pages": len(pages),
        "broken_links": broken_links,
        "orphan_pages": orphan_pages,
        "index_missing": index_missing,
        "manifest_missing": manifest_missing,
        "missing_required": missing_required,
        "quality": run_quality(vault, [p.relative_to(vault).as_posix() for p in pages if is_retrieval_page(p)], strict=False),
    }
    report_path = vault / "swarmvault" / "exports" / "lint-report.json"
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return report


def doctor(vault: Path) -> dict:
    vault = vault.expanduser()
    config = load_vault_config(vault)
    config_vault_root = config.get("vault_root") if isinstance(config, dict) else None
    config_vault_root_matches = None
    if config_vault_root:
        try:
            config_vault_root_matches = Path(config_vault_root).expanduser().resolve() == vault.resolve()
        except OSError:
            config_vault_root_matches = False
    obsidian_vaults = find_obsidian_vaults(vault)
    nested_obsidian_vaults = [path for path in obsidian_vaults if path != vault]
    warnings = []
    if not (vault / ".obsidian").is_dir() and nested_obsidian_vaults:
        warnings.append("selected path is not an Obsidian vault, but nested Obsidian vaults were found; use the folder containing .obsidian as --vault")
    if config_vault_root_matches is False:
        warnings.append("ghpf.config.json vault_root does not match the selected --vault path")
    return {
        "vault": str(vault),
        "resolved_vault": str(vault.resolve()),
        "has_obsidian_config": (vault / ".obsidian").is_dir(),
        "nested_obsidian_vaults": [str(path) for path in nested_obsidian_vaults],
        "has_raw": (vault / "_raw").exists(),
        "has_karpathy_raw": (vault / "raw").exists(),
        "has_graphify_raw": (vault / "raw" / "graphify_articles").exists(),
        "has_wiki": (vault / "wiki").exists(),
        "has_graph_imports": (vault / "graph_imports").exists(),
        "has_schema": (vault / "schema" / "AGENTS.md").exists(),
        "has_sidecar": (vault / "swarmvault").exists(),
        "has_cache": (vault / "swarmvault" / "cache").exists(),
        "markdown_pages": len(markdown_files(vault)),
        "reference_markdown_pages": len(reference_markdown_files(vault)),
        "has_config": (vault / "ghpf.config.json").exists(),
        "config_vault_root": config_vault_root,
        "config_vault_root_matches": config_vault_root_matches,
        "has_manifest": manifest_path(vault).exists(),
        "warnings": warnings,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    sub = parser.add_subparsers(dest="command", required=True)

    graph_p = sub.add_parser("graph")
    graph_p.add_argument("--vault", default=".")

    extract_p = sub.add_parser("extract")
    extract_p.add_argument("--vault", default=".")
    extract_p.add_argument("--ingest", action="store_true", help="Ingest extracted Markdown into wiki notes after extraction.")
    extract_p.add_argument("sources", nargs="+", help="PDF files, HTML files, web URLs, YouTube URLs, or text files.")

    ingest_p = sub.add_parser("ingest")
    ingest_p.add_argument("--vault", default=".")
    ingest_p.add_argument("sources", nargs="*", help="Files or URLs to ingest. Defaults to raw/ and _raw/.")
    ingest_p.add_argument("--move", action="store_true", help="Delete original external source after copying to raw/.")

    lint_p = sub.add_parser("lint")
    lint_p.add_argument("--vault", default=".")

    file_back_p = sub.add_parser("file-back")
    file_back_p.add_argument("--vault", default=".")
    file_back_p.add_argument("--title", required=True)
    file_back_p.add_argument("--body", required=True)
    file_back_p.add_argument("--folder", default="wiki/syntheses")

    context_p = sub.add_parser("context")
    context_p.add_argument("--vault", default=".")
    context_p.add_argument("--query", required=True)
    context_p.add_argument("--limit", type=int, default=8)

    task_p = sub.add_parser("task")
    task_sub = task_p.add_subparsers(dest="task_command", required=True)
    for name in ("start", "update", "finish"):
        p = task_sub.add_parser(name)
        p.add_argument("--vault", default=".")
        p.add_argument("--title", required=True)
        p.add_argument("--note", default="")
        p.add_argument("--target", default="")

    doctor_p = sub.add_parser("doctor")
    doctor_p.add_argument("--vault", default=".")

    quality_p = sub.add_parser("quality")
    quality_p.add_argument("--vault", default=".")
    quality_p.add_argument("--strict", action="store_true")
    quality_p.add_argument("--coverage", type=float)
    quality_p.add_argument("--consistency", type=float)
    quality_p.add_argument("--suggestions", type=float)
    quality_p.add_argument("pages", nargs="*", help="Wiki page paths relative to the vault. Defaults to all wiki pages.")

    diff_p = sub.add_parser("diff")
    diff_p.add_argument("--vault", default=".")
    diff_p.add_argument("existing", help="Existing Markdown page path, absolute or vault-relative.")
    diff_p.add_argument("new", help="New Markdown page path, absolute or vault-relative.")

    state_p = sub.add_parser("state")
    state_p.add_argument("--vault", default=".")
    state_p.add_argument("action", choices=["init", "show", "check", "complete"])
    state_p.add_argument("step", nargs="?", choices=PIPELINE_STEPS)

    link_audit_p = sub.add_parser("link-audit")
    link_audit_p.add_argument("--vault", default=".")

    link_strengthen_p = sub.add_parser("link-strengthen")
    link_strengthen_p.add_argument("--vault", default=".")
    link_strengthen_p.add_argument("--page", required=True, help="Wiki page path relative to the vault.")
    link_strengthen_p.add_argument("--max-links", type=int, default=5)
    link_strengthen_p.add_argument("--backlink", action="store_true")

    capabilities_p = sub.add_parser("capabilities")
    capabilities_p.add_argument("--vault")

    graphify_p = sub.add_parser("graphify-import")
    graphify_p.add_argument("--vault", default=".")
    graphify_p.add_argument("--graph", required=True, help="Graphify graph.json path.")
    graphify_p.add_argument("--run-id")
    graphify_p.add_argument("--report", help="Optional GRAPH_REPORT.md path.")
    graphify_p.add_argument("--html", help="Optional graph.html path.")
    graphify_p.add_argument("--max-links", type=int, default=20)

    cache_clean_p = sub.add_parser("cache-clean")
    cache_clean_p.add_argument("--vault", default=".")
    cache_clean_p.add_argument("--max-age-days", type=int, default=30)
    cache_clean_p.add_argument("--keep-latest", type=int, default=10)
    cache_clean_p.add_argument("--include-exports", action="store_true")
    cache_clean_p.add_argument("--dry-run", action="store_true")

    maintenance_p = sub.add_parser("maintenance")
    maintenance_p.add_argument("--vault", default=".")
    maintenance_p.add_argument("--threshold", type=int, default=20, help="Run or recommend Graphify after this many new source notes since the last Graphify import.")
    maintenance_p.add_argument("--auto-graphify", action="store_true", help="Import a Graphify graph when the threshold is met and --graphify-graph is provided.")
    maintenance_p.add_argument("--force-graphify", action="store_true", help="Treat Graphify as needed regardless of the threshold.")
    maintenance_p.add_argument("--graphify-graph", help="Existing Graphify graph.json to import when Graphify is needed.")
    maintenance_p.add_argument("--graphify-run-id")

    card_p = sub.add_parser("card")
    card_p.add_argument("--vault", default=".")
    card_p.add_argument("--type", choices=["paper", "experiment", "strategy"], required=True)
    card_p.add_argument("--all-sources", action="store_true")
    card_p.add_argument("pages", nargs="*", help="Wiki source pages to turn into structured cards.")

    index_p = sub.add_parser("index")
    index_p.add_argument("--vault", default=".")

    search_p = sub.add_parser("search")
    search_p.add_argument("--vault", default=".")
    search_p.add_argument("--query", required=True)
    search_p.add_argument("--limit", type=int, default=10)

    insight_p = sub.add_parser("insight")
    insight_p.add_argument("--vault", default=".")
    insight_p.add_argument("--type", choices=["paper", "experiment", "strategy"], required=True)
    insight_p.add_argument("--query", required=True)
    insight_p.add_argument("--limit", type=int, default=8)

    evaluate_p = sub.add_parser("evaluate")
    evaluate_p.add_argument("--vault", default=".")
    evaluate_p.add_argument("--type", choices=["paper", "experiment", "strategy"], required=True)
    evaluate_p.add_argument("--target", required=True)

    figure_card_p = sub.add_parser("figure-card")
    figure_card_p.add_argument("--vault", default=".")
    figure_card_p.add_argument("--domain", choices=["auto", "irrigation", "trading", "generic"], default="auto")
    figure_card_p.add_argument("--all-sources", action="store_true")
    figure_card_p.add_argument("pages", nargs="*", help="Wiki source pages to turn into figure cards.")

    figure_insight_p = sub.add_parser("figure-insight")
    figure_insight_p.add_argument("--vault", default=".")
    figure_insight_p.add_argument("--domain", choices=["auto", "irrigation", "trading", "generic"], default="auto")
    figure_insight_p.add_argument("--query", required=True)
    figure_insight_p.add_argument("--data-schema")
    figure_insight_p.add_argument("--limit", type=int, default=8)

    figure_export_p = sub.add_parser("figure-export")
    figure_export_p.add_argument("--vault", default=".")
    figure_export_p.add_argument("--design")
    figure_export_p.add_argument("--domain", choices=["auto", "irrigation", "trading", "generic"], default="auto")
    figure_export_p.add_argument("--name", default="Figure_1")
    figure_export_p.add_argument("--data", help="Optional CSV path with real data.")
    figure_export_p.add_argument("--run", action="store_true", help="Run generated Matplotlib script and export PDF/SVG/PNG.")

    video_frames_p = sub.add_parser("video-frames")
    video_frames_p.add_argument("--vault", default=".")
    video_frames_p.add_argument("--source", required=True, help="YouTube URL, remote image/video URL, local image, or local video.")
    video_frames_p.add_argument("--every-seconds", type=float, default=30.0, help="Frame interval for video sources.")
    video_frames_p.add_argument("--max-frames", type=int, default=12)
    video_frames_p.add_argument("--ocr", action="store_true", help="Run OCR when tesseract or pytesseract is available.")
    video_frames_p.add_argument("--ocr-languages", default="eng+kor")
    video_frames_p.add_argument("--ingest", action="store_true", help="Ingest the generated frame-analysis Markdown into the wiki.")
    video_frames_p.add_argument("--figure-card", action="store_true", help="Create a figure card after ingesting the generated source note.")

    args = parser.parse_args()
    if args.command == "extract":
        result = extract_sources(Path(args.vault), args.sources)
        if args.ingest:
            paths = [item["path"] for item in result["extracted"]]
            result["ingest"] = ingest_sources(Path(args.vault), paths)
        print(json.dumps(result, ensure_ascii=False, indent=2))
    elif args.command == "ingest":
        print(json.dumps(ingest_sources(Path(args.vault), args.sources, move=args.move), ensure_ascii=False, indent=2))
    elif args.command == "lint":
        report = lint_wiki(Path(args.vault))
        print(json.dumps(report, ensure_ascii=False, indent=2))
        if not report["ok"]:
            return 1
    elif args.command == "file-back":
        print(file_back(Path(args.vault), args.title, args.body, folder=args.folder))
    elif args.command == "graph":
        graph = build_graph(Path(args.vault).expanduser())
        print(json.dumps({"node_count": graph["node_count"], "edge_count": graph["edge_count"]}, indent=2))
    elif args.command == "context":
        path = build_context(Path(args.vault).expanduser(), args.query, limit=args.limit)
        print(path)
    elif args.command == "task":
        path = record_task(Path(args.vault).expanduser(), args.task_command, args.title, args.note, args.target)
        print(path)
    elif args.command == "doctor":
        print(json.dumps(doctor(Path(args.vault).expanduser()), ensure_ascii=False, indent=2))
    elif args.command == "quality":
        report = run_quality(
            Path(args.vault),
            args.pages,
            strict=args.strict,
            consistency=args.consistency,
            suggestions=args.suggestions,
            coverage=args.coverage,
        )
        print(json.dumps(report, ensure_ascii=False, indent=2))
        if not report["ok"]:
            return 1
    elif args.command == "diff":
        vault = Path(args.vault)
        report = section_diff(vault_relative_path(vault, args.existing), vault_relative_path(vault, args.new))
        out = vault.expanduser() / "swarmvault" / "exports" / "section-diff.json"
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        print(json.dumps(report, ensure_ascii=False, indent=2))
    elif args.command == "state":
        print(json.dumps(state_action(Path(args.vault), args.action, args.step), ensure_ascii=False, indent=2))
    elif args.command == "link-audit":
        print(json.dumps(link_audit(Path(args.vault)), ensure_ascii=False, indent=2))
    elif args.command == "link-strengthen":
        print(json.dumps(link_strengthen(Path(args.vault), args.page, max_links=args.max_links, backlink=args.backlink), ensure_ascii=False, indent=2))
    elif args.command == "capabilities":
        vault = Path(args.vault).expanduser() if args.vault else None
        print(json.dumps(capabilities(vault), ensure_ascii=False, indent=2))
    elif args.command == "graphify-import":
        graph_path = Path(args.graph).expanduser()
        report = Path(args.report).expanduser() if args.report else graph_path.parent / "GRAPH_REPORT.md"
        html_path = Path(args.html).expanduser() if args.html else graph_path.parent / "graph.html"
        report = report if report.exists() else None
        html_path = html_path if html_path.exists() else None
        result = graphify_import(Path(args.vault), graph_path, run_id=args.run_id, report=report, html_file=html_path, max_links=args.max_links)
        print(json.dumps(result, ensure_ascii=False, indent=2))
    elif args.command == "cache-clean":
        print(
            json.dumps(
                cache_clean(
                    Path(args.vault),
                    max_age_days=args.max_age_days,
                    keep_latest=args.keep_latest,
                    dry_run=args.dry_run,
                    include_exports=args.include_exports,
                ),
                ensure_ascii=False,
                indent=2,
            )
        )
    elif args.command == "maintenance":
        graphify_graph = Path(args.graphify_graph).expanduser() if args.graphify_graph else None
        print(
            json.dumps(
                maintenance(
                    Path(args.vault),
                    threshold=args.threshold,
                    auto_graphify=args.auto_graphify,
                    force_graphify=args.force_graphify,
                    graphify_graph=graphify_graph,
                    graphify_run_id=args.graphify_run_id,
                ),
                ensure_ascii=False,
                indent=2,
            )
        )
    elif args.command == "card":
        print(json.dumps(card_command(Path(args.vault), args.type, args.pages, all_sources=args.all_sources), ensure_ascii=False, indent=2))
    elif args.command == "index":
        print(json.dumps(build_hybrid_index(Path(args.vault)), ensure_ascii=False, indent=2))
    elif args.command == "search":
        print(json.dumps(hybrid_search(Path(args.vault), args.query, limit=args.limit), ensure_ascii=False, indent=2))
    elif args.command == "insight":
        print(json.dumps(insight_command(Path(args.vault), args.type, args.query, limit=args.limit), ensure_ascii=False, indent=2))
    elif args.command == "evaluate":
        print(json.dumps(evaluate_command(Path(args.vault), args.type, args.target), ensure_ascii=False, indent=2))
    elif args.command == "figure-card":
        print(json.dumps(figure_card_command(Path(args.vault), args.pages, domain=args.domain, all_sources=args.all_sources), ensure_ascii=False, indent=2))
    elif args.command == "figure-insight":
        print(json.dumps(figure_insight_command(Path(args.vault), args.query, domain=args.domain, data_schema=args.data_schema, limit=args.limit), ensure_ascii=False, indent=2))
    elif args.command == "figure-export":
        print(json.dumps(figure_export_command(Path(args.vault), args.design, domain=args.domain, name=args.name, data=args.data, run=args.run), ensure_ascii=False, indent=2))
    elif args.command == "video-frames":
        print(
            json.dumps(
                video_frames_command(
                    Path(args.vault),
                    args.source,
                    every_seconds=args.every_seconds,
                    max_frames=args.max_frames,
                    ocr=args.ocr,
                    ocr_languages=args.ocr_languages,
                    ingest=args.ingest,
                    figure_card=args.figure_card,
                ),
                ensure_ascii=False,
                indent=2,
            )
        )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
